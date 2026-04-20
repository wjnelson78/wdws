#!/usr/bin/env python3
"""
OneDrive Sync Service for Athena Chat Database.

Connects to Microsoft Graph API using app-only (client credentials) auth,
syncs OneDrive/SharePoint files, extracts text content, generates embeddings,
and stores everything in the athena_chat PostgreSQL database.

Architecture:
  - Layer 1: Global index (chat.onedrive_files + chat.onedrive_file_chunks)
  - Layer 2: Space linking via chat.space_onedrive_links (separate concern)

Supports:
  - Full sync (crawl entire drive/folder tree)
  - Delta sync (Graph delta API for incremental changes)
  - Webhook-triggered re-sync on file changes
  - Selective file extension filtering
  - Multiple sync targets (different drives/folders)

Usage:
    python onedrive_sync.py                        # Sync all active targets
    python onedrive_sync.py --target 1             # Sync specific target ID
    python onedrive_sync.py --add-target user william@seattleseahawks.me /
    python onedrive_sync.py --add-target user william@seattleseahawks.me /Documents/Legal
    python onedrive_sync.py --list-targets
    python onedrive_sync.py --dry-run
    python onedrive_sync.py --daemon               # Run as continuous service
"""
import os
import sys
sys.path.insert(0, "/opt/wdws")
import json
import re
import uuid
import hashlib
import asyncio
import logging
import time
import signal
import argparse
import io
import tempfile
from pathlib import Path
from datetime import datetime, timedelta, timezone
from dataclasses import dataclass, field
from typing import List, Dict, Optional, Set, Tuple, Any

import httpx
import asyncpg

from embedding_service import (
    embed_texts_sync, embed_query_sync,
    _vec_literal as _embedding_vec_literal,
    EMBEDDING_DIMENSIONS, EMBEDDING_MODEL,
)
from contextual_retrieval import generate_context_sync, enrich_chunks

# ============================================================
# Configuration
# ============================================================

_env_file = Path("/opt/wdws/.env")
if _env_file.exists():
    for _line in _env_file.read_text().splitlines():
        _line = _line.strip()
        if _line and not _line.startswith("#") and "=" in _line:
            _k, _v = _line.split("=", 1)
            os.environ.setdefault(_k.strip(), _v.strip())

CHAT_DATABASE_URL = os.getenv("CHAT_DATABASE_URL", "")
if not CHAT_DATABASE_URL:
    raise RuntimeError("Missing CHAT_DATABASE_URL — check /opt/wdws/.env")

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
# OpenAI API key no longer required — embeddings use local BGE-M3 model

GRAPH_CLIENT_ID = os.getenv("GRAPH_CLIENT_ID", "")
GRAPH_CLIENT_SECRET = os.getenv("GRAPH_CLIENT_SECRET", "")
GRAPH_TENANT_ID = os.getenv("GRAPH_TENANT_ID", "")

if not all([GRAPH_CLIENT_ID, GRAPH_CLIENT_SECRET, GRAPH_TENANT_ID]):
    raise RuntimeError("Missing GRAPH_CLIENT_ID/SECRET/TENANT_ID — check /opt/wdws/.env")

GRAPH_BASE_URL = "https://graph.microsoft.com/v1.0"
TOKEN_URL = f"https://login.microsoftonline.com/{GRAPH_TENANT_ID}/oauth2/v2.0/token"

# EMBEDDING_MODEL and EMBEDDING_DIMENSIONS imported from embedding_service
# (BGE-M3 local model, 1024 dimensions)
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
EMBEDDING_BATCH_SIZE = 20          # chunks per OpenAI API call
MAX_FILE_DOWNLOAD_SIZE = 100_000_000   # 100MB
SYNC_INTERVAL_SECONDS = 300        # 5 minutes in daemon mode
WEBHOOK_NOTIFICATION_URL = os.getenv(
    "ONEDRIVE_WEBHOOK_URL", "https://athena-ai.12432.net/webhooks/graph"
)
WEBHOOK_EXPIRY_HOURS = 48          # Graph max is 4230 min (~70.5 hours)

# Supported file extensions for text extraction
SUPPORTED_EXTENSIONS = {
    # Office
    'docx', 'doc', 'xlsx', 'xls', 'pptx', 'ppt',
    # PDF
    'pdf',
    # Text
    'txt', 'md', 'csv', 'tsv', 'json', 'xml', 'html', 'htm',
    'log', 'yaml', 'yml', 'ini', 'cfg', 'conf',
    # Code (useful for dev teams)
    'py', 'js', 'ts', 'sql', 'sh', 'ps1', 'psm1',
    'cs', 'java', 'cpp', 'c', 'h', 'rb', 'go', 'rs',
}

# ============================================================
# Logging
# ============================================================

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s [%(name)s] %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("onedrive_sync")


# ============================================================
# Data classes
# ============================================================

@dataclass
class SyncTarget:
    id: int
    drive_type: str
    drive_owner: str
    drive_id: Optional[str]
    folder_path: str
    folder_id: Optional[str]
    is_active: bool
    sync_mode: str
    file_extensions: Optional[List[str]]
    max_file_size: int
    recursive: bool
    delta_link: Optional[str]
    label: Optional[str]


@dataclass
class DriveItem:
    """Represents a file from the Graph API."""
    item_id: str
    drive_id: str
    name: str
    path: str
    size: int
    mime_type: str
    extension: str
    created_by: Optional[str]
    created_by_email: Optional[str]
    modified_by: Optional[str]
    modified_by_email: Optional[str]
    web_url: Optional[str]
    created_at: Optional[datetime]
    modified_at: Optional[datetime]
    etag: Optional[str]
    ctag: Optional[str]
    parent_item_id: Optional[str]
    is_deleted: bool = False
    download_url: Optional[str] = None


@dataclass
class SyncStats:
    files_discovered: int = 0
    files_new: int = 0
    files_updated: int = 0
    files_deleted: int = 0
    files_skipped: int = 0
    files_failed: int = 0
    chunks_created: int = 0
    embedding_tokens: int = 0


# ============================================================
# Text chunker (same logic as email_sync.py)
# ============================================================

class TextChunker:
    def __init__(self, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = ["\n\n", "\n", ". ", " ", ""]

    def split(self, text: str) -> List[str]:
        if not text:
            return []
        return self._split_text(text, self.separators)

    def _split_text(self, text: str, separators: list) -> list:
        final_chunks = []
        separator = separators[-1]
        new_separators = []
        for i, sep in enumerate(separators):
            if sep == "" or sep in text:
                separator = sep
                new_separators = separators[i + 1:]
                break
        splits = text.split(separator) if separator else list(text)
        good_splits = []
        for s in splits:
            if len(s) < self.chunk_size:
                good_splits.append(s)
            else:
                if good_splits:
                    merged = self._merge_splits(good_splits, separator)
                    final_chunks.extend(merged)
                    good_splits = []
                if new_separators:
                    final_chunks.extend(self._split_text(s, new_separators))
                else:
                    final_chunks.append(s)
        if good_splits:
            merged = self._merge_splits(good_splits, separator)
            final_chunks.extend(merged)
        return final_chunks

    def _merge_splits(self, splits: list, separator: str) -> list:
        docs = []
        current = []
        total = 0
        for s in splits:
            s_len = len(s)
            if total + s_len + (len(separator) if current else 0) > self.chunk_size:
                if current:
                    docs.append(separator.join(current))
                    while total > self.chunk_overlap and len(current) > 1:
                        total -= len(current[0]) + len(separator)
                        current.pop(0)
                current = [s]
                total = s_len
            else:
                current.append(s)
                total += s_len + (len(separator) if len(current) > 1 else 0)
        if current:
            docs.append(separator.join(current))
        return docs


# ============================================================
# Content extractors
# ============================================================

def extract_text_from_docx(content: bytes) -> str:
    """Extract text from a .docx file."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        return "\n\n".join(paragraphs)
    except Exception as e:
        log.warning(f"  DOCX extraction failed: {e}")
        return ""


def extract_text_from_pdf(content: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(content))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i+1}]\n{text}")
        return "\n\n".join(pages)
    except Exception as e:
        log.warning(f"  PDF extraction failed: {e}")
        return ""


def extract_text_from_xlsx(content: bytes) -> str:
    """Extract text from an Excel file."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                cells = [str(c).strip() for c in row if c is not None]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except Exception as e:
        log.warning(f"  XLSX extraction failed: {e}")
        return ""


def extract_text_from_pptx(content: bytes) -> str:
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            texts.append(para.text)
            if texts:
                slides.append(f"[Slide {i+1}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as e:
        log.warning(f"  PPTX extraction failed: {e}")
        return ""


def extract_text_from_csv(content: bytes) -> str:
    """Extract text from a CSV/TSV file."""
    try:
        import csv
        text = content.decode("utf-8", errors="ignore")
        reader = csv.reader(io.StringIO(text))
        rows = []
        for row in reader:
            cells = [c.strip() for c in row if c.strip()]
            if cells:
                rows.append(" | ".join(cells))
        return "\n".join(rows[:5000])  # Cap at 5000 rows
    except Exception as e:
        log.warning(f"  CSV extraction failed: {e}")
        return ""


def extract_text_from_plain(content: bytes) -> str:
    """Extract text from plain text files."""
    try:
        return content.decode("utf-8", errors="ignore")[:500000]
    except Exception:
        return ""


def extract_text(content: bytes, extension: str, filename: str) -> Tuple[str, str]:
    """
    Extract text from file content based on extension.
    Returns (text, extraction_method).
    """
    ext = extension.lower().lstrip('.')
    
    if ext == 'docx':
        return extract_text_from_docx(content), 'docx'
    elif ext == 'pdf':
        return extract_text_from_pdf(content), 'pdf'
    elif ext in ('xlsx', 'xls'):
        return extract_text_from_xlsx(content), 'xlsx'
    elif ext in ('pptx', 'ppt'):
        return extract_text_from_pptx(content), 'pptx'
    elif ext in ('csv', 'tsv'):
        return extract_text_from_csv(content), 'csv'
    elif ext in ('txt', 'md', 'log', 'json', 'xml', 'html', 'htm',
                 'yaml', 'yml', 'ini', 'cfg', 'conf',
                 'py', 'js', 'ts', 'sql', 'sh', 'ps1', 'psm1',
                 'cs', 'java', 'cpp', 'c', 'h', 'rb', 'go', 'rs'):
        return extract_text_from_plain(content), 'plain'
    else:
        return "", 'unsupported'


# ============================================================
# Embedding client
# ============================================================

class EmbeddingClient:
    """Local BGE-M3 embedding client with batching (delegates to embedding_service)."""

    def __init__(self, **kwargs):
        self.total_tokens = 0
        self.total_requests = 0

    async def embed_batch(self, texts: List[str]) -> List[List[float]]:
        if not texts:
            return []
        result = embed_texts_sync(texts)
        self.total_requests += 1
        return result

    async def close(self):
        pass  # No HTTP client to close


# ============================================================
# Graph API client (extended for OneDrive)
# ============================================================

class GraphClient:
    """Microsoft Graph API client using app-only (client credentials) auth."""

    def __init__(self):
        self.client = httpx.AsyncClient(timeout=60.0)
        self.access_token: Optional[str] = None
        self.token_expires_at: float = 0

    async def authenticate(self) -> bool:
        try:
            resp = await self.client.post(
                TOKEN_URL,
                data={
                    "client_id": GRAPH_CLIENT_ID,
                    "client_secret": GRAPH_CLIENT_SECRET,
                    "scope": "https://graph.microsoft.com/.default",
                    "grant_type": "client_credentials",
                },
            )
            resp.raise_for_status()
            data = resp.json()
            self.access_token = data["access_token"]
            self.token_expires_at = time.time() + data.get("expires_in", 3600) - 300
            log.info("✓ Graph API authenticated (app-only)")
            return True
        except Exception as e:
            log.error(f"✗ Graph API auth failed: {e}")
            return False

    async def _ensure_token(self):
        if time.time() >= self.token_expires_at:
            await self.authenticate()

    async def _get(self, url: str, params: dict = None) -> Optional[dict]:
        for attempt in range(5):
            await self._ensure_token()
            try:
                resp = await self.client.get(
                    url,
                    headers={"Authorization": f"Bearer {self.access_token}"},
                    params=params,
                )
                resp.raise_for_status()
                return resp.json()
            except httpx.HTTPStatusError as e:
                if e.response.status_code == 429 and attempt < 4:
                    retry_after = int(e.response.headers.get("Retry-After", "10"))
                    log.warning(f"  Rate limited (attempt {attempt+1}/5), waiting {retry_after}s...")
                    await asyncio.sleep(retry_after)
                    continue
                if e.response.status_code == 404:
                    log.warning(f"  Not found: {url}")
                    return None
                log.error(f"  Graph API error {e.response.status_code}: {e.response.text[:200]}")
                return None
            except Exception as e:
                log.error(f"  Graph API request failed: {e}")
                return None
        return None

    async def _get_bytes(self, url: str) -> Optional[bytes]:
        """Download file content as bytes."""
        for attempt in range(5):
            await self._ensure_token()
            try:
                resp = await self.client.get(
                    url,
                    headers={"Authorization": f"Bearer {self.access_token}"},
                    follow_redirects=True,
                )
                resp.raise_for_status()
                return resp.content
            except httpx.HTTPStatusError as e:
                if e.response.status_code == 429 and attempt < 4:
                    retry_after = int(e.response.headers.get("Retry-After", "10"))
                    log.warning(f"  Rate limited on download (attempt {attempt+1}/5), waiting {retry_after}s...")
                    await asyncio.sleep(retry_after)
                    continue
                log.error(f"  Download error {e.response.status_code}")
                return None
            except Exception as e:
                log.error(f"  Download failed: {e}")
                return None
        return None

    # ── Drive resolution ──

    async def resolve_drive_id(self, target: SyncTarget) -> Optional[str]:
        """Resolve a sync target to a Graph drive ID."""
        if target.drive_id:
            return target.drive_id

        if target.drive_type == 'user':
            data = await self._get(f"{GRAPH_BASE_URL}/users/{target.drive_owner}/drive")
            if data:
                return data.get("id")
        elif target.drive_type == 'site':
            # target.drive_owner is site URL or site-id
            data = await self._get(f"{GRAPH_BASE_URL}/sites/{target.drive_owner}/drive")
            if data:
                return data.get("id")
        elif target.drive_type == 'group':
            data = await self._get(f"{GRAPH_BASE_URL}/groups/{target.drive_owner}/drive")
            if data:
                return data.get("id")
        return None

    async def resolve_folder_id(self, drive_id: str, folder_path: str) -> Optional[str]:
        """Resolve a folder path to an item ID."""
        if folder_path == '/' or not folder_path:
            return 'root'
        
        # Strip leading/trailing slashes for Graph API path format
        clean_path = folder_path.strip('/')
        data = await self._get(
            f"{GRAPH_BASE_URL}/drives/{drive_id}/root:/{clean_path}"
        )
        if data:
            return data.get("id")
        return None

    # ── File listing ──

    async def list_children(self, drive_id: str, item_id: str = "root",
                            select: str = None) -> List[dict]:
        """List immediate children of a drive item."""
        url = f"{GRAPH_BASE_URL}/drives/{drive_id}/items/{item_id}/children"
        params = {"$top": 200}
        if select:
            params["$select"] = select
        
        all_items = []
        while url:
            data = await self._get(url, params if '@odata.nextLink' not in (url or '') else None)
            if not data:
                break
            all_items.extend(data.get("value", []))
            url = data.get("@odata.nextLink")
            params = None  # nextLink has params built in
        return all_items

    async def crawl_folder(self, drive_id: str, folder_id: str = "root",
                           path_prefix: str = "/", recursive: bool = True,
                           extensions: Optional[Set[str]] = None,
                           max_size: int = MAX_FILE_DOWNLOAD_SIZE) -> List[DriveItem]:
        """Recursively crawl a folder and return all matching file items."""
        items = await self.list_children(drive_id, folder_id)
        results: List[DriveItem] = []
        
        for item in items:
            name = item.get("name", "")
            item_id = item.get("id", "")
            
            if "folder" in item:
                # It's a folder — recurse
                if recursive:
                    sub_path = f"{path_prefix}{name}/"
                    sub_items = await self.crawl_folder(
                        drive_id, item_id, sub_path, recursive, extensions, max_size
                    )
                    results.extend(sub_items)
                continue
            
            if "file" not in item:
                continue  # Skip non-file items (notebooks, etc.)
            
            # It's a file
            ext = name.rsplit('.', 1)[-1].lower() if '.' in name else ''
            size = item.get("size", 0)
            
            # Filter by extension
            if extensions and ext not in extensions:
                continue
            
            # Filter by size
            if size > max_size:
                log.info(f"  ⊘ Skipping {name} ({size/1024/1024:.1f}MB > {max_size/1024/1024:.0f}MB limit)")
                continue
            
            # Parse metadata
            created_by_info = item.get("createdBy", {}).get("user", {})
            modified_by_info = item.get("lastModifiedBy", {}).get("user", {})
            parent_ref = item.get("parentReference", {})
            
            drive_item = DriveItem(
                item_id=item_id,
                drive_id=drive_id,
                name=name,
                path=f"{path_prefix}{name}",
                size=size,
                mime_type=item.get("file", {}).get("mimeType", ""),
                extension=ext,
                created_by=created_by_info.get("displayName"),
                created_by_email=created_by_info.get("email"),
                modified_by=modified_by_info.get("displayName"),
                modified_by_email=modified_by_info.get("email"),
                web_url=item.get("webUrl"),
                created_at=_parse_graph_date(item.get("createdDateTime")),
                modified_at=_parse_graph_date(item.get("lastModifiedDateTime")),
                etag=item.get("eTag"),
                ctag=item.get("cTag"),
                parent_item_id=parent_ref.get("id"),
                download_url=item.get("@microsoft.graph.downloadUrl"),
            )
            results.append(drive_item)
        
        return results

    async def get_delta(self, drive_id: str, folder_id: str = "root",
                        delta_link: Optional[str] = None) -> Tuple[List[dict], Optional[str]]:
        """
        Use delta API for incremental sync.
        Returns (changed_items, new_delta_link).
        """
        if delta_link:
            url = delta_link
        else:
            if folder_id == "root":
                url = f"{GRAPH_BASE_URL}/drives/{drive_id}/root/delta"
            else:
                url = f"{GRAPH_BASE_URL}/drives/{drive_id}/items/{folder_id}/delta"
        
        all_items = []
        new_delta_link = None
        
        while url:
            data = await self._get(url)
            if not data:
                break
            all_items.extend(data.get("value", []))
            url = data.get("@odata.nextLink")
            if "@odata.deltaLink" in data:
                new_delta_link = data["@odata.deltaLink"]
        
        return all_items, new_delta_link

    async def download_file(self, drive_id: str, item_id: str) -> Optional[bytes]:
        """Download file content."""
        url = f"{GRAPH_BASE_URL}/drives/{drive_id}/items/{item_id}/content"
        return await self._get_bytes(url)

    # ── Webhook subscriptions ──

    async def create_subscription(self, resource: str, notification_url: str,
                                   expiration: datetime) -> Optional[dict]:
        """Create a Graph API change notification subscription."""
        await self._ensure_token()
        # Format expiration as ISO 8601 UTC (Graph requires trailing Z, no +00:00)
        exp_str = expiration.strftime("%Y-%m-%dT%H:%M:%S.000Z")
        payload = {
            "changeType": "updated",
            "notificationUrl": notification_url,
            "resource": resource,
            "expirationDateTime": exp_str,
            "clientState": "wdws-onedrive-sync",
        }
        log.info(f"  Subscription payload: resource={resource} url={notification_url} expires={exp_str}")
        try:
            resp = await self.client.post(
                f"{GRAPH_BASE_URL}/subscriptions",
                headers={
                    "Authorization": f"Bearer {self.access_token}",
                    "Content-Type": "application/json",
                },
                json=payload,
            )
            if resp.status_code >= 400:
                log.error(f"  Subscription API error {resp.status_code}: {resp.text}")
            resp.raise_for_status()
            return resp.json()
        except Exception as e:
            log.error(f"  Subscription creation failed: {e}")
            return None

    async def renew_subscription(self, subscription_id: str, expiration: datetime) -> bool:
        """Renew an existing subscription."""
        await self._ensure_token()
        try:
            resp = await self.client.patch(
                f"{GRAPH_BASE_URL}/subscriptions/{subscription_id}",
                headers={
                    "Authorization": f"Bearer {self.access_token}",
                    "Content-Type": "application/json",
                },
                json={
                    "expirationDateTime": expiration.isoformat() + "Z",
                },
            )
            resp.raise_for_status()
            return True
        except Exception as e:
            log.error(f"  Subscription renewal failed: {e}")
            return False

    async def close(self):
        await self.client.aclose()


# ============================================================
# Helpers
# ============================================================

def _parse_graph_date(date_str: Optional[str]) -> Optional[datetime]:
    """Parse Graph API datetime string."""
    if not date_str:
        return None
    try:
        return datetime.fromisoformat(date_str.replace("Z", "+00:00"))
    except Exception:
        return None


def _vec_literal(embedding: List[float]) -> str:
    """Convert embedding list to PostgreSQL halfvec literal."""
    return _embedding_vec_literal(embedding)


# ============================================================
# Database operations
# ============================================================

async def load_sync_targets(pool: asyncpg.Pool, target_id: Optional[int] = None) -> List[SyncTarget]:
    """Load active sync targets from DB."""
    sql = """
        SELECT id, drive_type, drive_owner, drive_id, folder_path, folder_id,
               is_active, sync_mode, file_extensions, max_file_size, recursive,
               delta_link, label
        FROM chat.onedrive_sync_targets
        WHERE is_active = true
    """
    params = []
    if target_id:
        sql += " AND id = $1"
        params.append(target_id)
    sql += " ORDER BY id"
    
    rows = await pool.fetch(sql, *params)
    return [SyncTarget(
        id=r["id"], drive_type=r["drive_type"], drive_owner=r["drive_owner"],
        drive_id=r["drive_id"], folder_path=r["folder_path"], folder_id=r["folder_id"],
        is_active=r["is_active"], sync_mode=r["sync_mode"],
        file_extensions=r["file_extensions"], max_file_size=r["max_file_size"],
        recursive=r["recursive"], delta_link=r["delta_link"], label=r["label"],
    ) for r in rows]


async def add_sync_target(pool: asyncpg.Pool, drive_type: str, drive_owner: str,
                          folder_path: str = "/", label: str = None,
                          extensions: List[str] = None, sync_mode: str = "full") -> int:
    """Add a new sync target. Returns the target ID."""
    return await pool.fetchval("""
        INSERT INTO chat.onedrive_sync_targets
            (drive_type, drive_owner, folder_path, label, file_extensions, sync_mode)
        VALUES ($1, $2, $3, $4, $5, $6)
        ON CONFLICT (drive_type, drive_owner, folder_path) DO UPDATE
            SET is_active = true, label = COALESCE(EXCLUDED.label, chat.onedrive_sync_targets.label)
        RETURNING id
    """, drive_type, drive_owner, folder_path, label, extensions, sync_mode)


async def get_existing_files(pool: asyncpg.Pool, sync_target_id: int) -> Dict[str, dict]:
    """Get all existing files for a sync target, keyed by (drive_id, item_id)."""
    rows = await pool.fetch("""
        SELECT id, drive_id, item_id, ctag, status
        FROM chat.onedrive_files
        WHERE sync_target_id = $1 AND status != 'deleted'
    """, sync_target_id)
    return {f"{r['drive_id']}:{r['item_id']}": dict(r) for r in rows}


async def upsert_file_record(pool: asyncpg.Pool, target: SyncTarget,
                              item: DriveItem, status: str = "pending") -> uuid.UUID:
    """Insert or update a file record. Returns the file UUID."""
    file_id = await pool.fetchval("""
        INSERT INTO chat.onedrive_files
            (sync_target_id, drive_id, item_id, parent_item_id,
             file_name, file_path, file_extension, mime_type, size_bytes,
             created_by, created_by_email, last_modified_by, last_modified_by_email,
             web_url, graph_created_at, graph_modified_at, etag, ctag,
             status, last_synced_at)
        VALUES ($1, $2, $3, $4, $5, $6, $7, $8, $9, $10, $11, $12, $13,
                $14, $15, $16, $17, $18, $19, now())
        ON CONFLICT (drive_id, item_id) DO UPDATE SET
            file_name = EXCLUDED.file_name,
            file_path = EXCLUDED.file_path,
            size_bytes = EXCLUDED.size_bytes,
            last_modified_by = EXCLUDED.last_modified_by,
            last_modified_by_email = EXCLUDED.last_modified_by_email,
            web_url = EXCLUDED.web_url,
            graph_modified_at = EXCLUDED.graph_modified_at,
            etag = EXCLUDED.etag,
            ctag = EXCLUDED.ctag,
            status = CASE 
                WHEN chat.onedrive_files.ctag != EXCLUDED.ctag THEN 'pending'
                ELSE chat.onedrive_files.status
            END,
            last_synced_at = now()
        RETURNING id
    """, target.id, item.drive_id, item.item_id, item.parent_item_id,
        item.name, item.path, item.extension, item.mime_type, item.size,
        item.created_by, item.created_by_email, item.modified_by, item.modified_by_email,
        item.web_url, item.created_at, item.modified_at, item.etag, item.ctag,
        status)
    return file_id


async def write_file_content(pool: asyncpg.Pool, file_id: uuid.UUID,
                              content_text: str, extraction_method: str,
                              chunks: List[str], embeddings: List[List[float]]):
    """Write extracted content and embeddings for a file."""
    async with pool.acquire() as conn:
        async with conn.transaction():
            # Update file record with content
            await conn.execute("""
                UPDATE chat.onedrive_files SET
                    content_text = $2,
                    extraction_method = $3,
                    status = 'indexed',
                    indexed_at = now(),
                    processing_error = NULL
                WHERE id = $1
            """, file_id, content_text[:500000])
            
            # Delete old chunks (for re-index)
            await conn.execute(
                "DELETE FROM chat.onedrive_file_chunks WHERE file_id = $1",
                file_id
            )
            
            # Insert chunks with embeddings
            for i, (chunk_text, emb) in enumerate(zip(chunks, embeddings)):
                await conn.execute("""
                    INSERT INTO chat.onedrive_file_chunks
                        (file_id, chunk_index, chunk_text, start_offset, end_offset,
                         token_count, embedding)
                    VALUES ($1, $2, $3, $4, $5, $6, $7::halfvec)
                """, file_id, i, chunk_text, 0, len(chunk_text),
                    len(chunk_text.split()), _vec_literal(emb))


async def mark_file_failed(pool: asyncpg.Pool, file_id: uuid.UUID, error: str):
    """Mark a file as failed."""
    await pool.execute("""
        UPDATE chat.onedrive_files SET
            status = 'failed',
            processing_error = $2
        WHERE id = $1
    """, file_id, error[:1000])


async def mark_file_skipped(pool: asyncpg.Pool, file_id: uuid.UUID, reason: str):
    """Mark a file as skipped (unsupported type, empty content, etc.)."""
    await pool.execute("""
        UPDATE chat.onedrive_files SET
            status = 'skipped',
            processing_error = $2
        WHERE id = $1
    """, file_id, reason[:500])


async def mark_files_deleted(pool: asyncpg.Pool, file_ids: List[uuid.UUID]):
    """Mark files as deleted (from delta sync)."""
    if not file_ids:
        return
    await pool.execute("""
        UPDATE chat.onedrive_files SET status = 'deleted'
        WHERE id = ANY($1::uuid[])
    """, file_ids)
    # Remove chunks for deleted files
    await pool.execute("""
        DELETE FROM chat.onedrive_file_chunks 
        WHERE file_id = ANY($1::uuid[])
    """, file_ids)


async def update_sync_target_stats(pool: asyncpg.Pool, target_id: int):
    """Update aggregate stats on the sync target."""
    await pool.execute("""
        UPDATE chat.onedrive_sync_targets SET
            total_files = (SELECT COUNT(*) FROM chat.onedrive_files 
                          WHERE sync_target_id = $1 AND status NOT IN ('deleted', 'skipped')),
            total_chunks = (SELECT COUNT(*) FROM chat.onedrive_file_chunks fc
                           JOIN chat.onedrive_files f ON f.id = fc.file_id
                           WHERE f.sync_target_id = $1),
            total_bytes = (SELECT COALESCE(SUM(size_bytes), 0) FROM chat.onedrive_files
                          WHERE sync_target_id = $1 AND status NOT IN ('deleted', 'skipped')),
            updated_at = now()
        WHERE id = $1
    """, target_id)


async def record_sync_run(pool: asyncpg.Pool, target_id: int,
                           trigger: str = "scheduled") -> int:
    """Create a sync run record. Returns run ID."""
    return await pool.fetchval("""
        INSERT INTO chat.onedrive_sync_runs (sync_target_id, trigger_type)
        VALUES ($1, $2) RETURNING id
    """, target_id, trigger)


async def complete_sync_run(pool: asyncpg.Pool, run_id: int, stats: SyncStats,
                             duration: float, delta_link: Optional[str] = None,
                             error: Optional[str] = None):
    """Finalize a sync run record."""
    status = "failed" if error else "completed"
    await pool.execute("""
        UPDATE chat.onedrive_sync_runs SET
            completed_at = now(),
            status = $2,
            files_discovered = $3,
            files_new = $4,
            files_updated = $5,
            files_deleted = $6,
            files_skipped = $7,
            files_failed = $8,
            chunks_created = $9,
            embedding_tokens = $10,
            duration_seconds = $11,
            error_message = $12,
            delta_link_after = $13
        WHERE id = $1
    """, run_id, status, stats.files_discovered, stats.files_new,
        stats.files_updated, stats.files_deleted, stats.files_skipped,
        stats.files_failed, stats.chunks_created, stats.embedding_tokens,
        duration, error, delta_link)


# ============================================================
# Core sync pipeline
# ============================================================

async def sync_target(pool: asyncpg.Pool, graph: GraphClient,
                      embedder: EmbeddingClient, chunker: TextChunker,
                      target: SyncTarget, dry_run: bool = False,
                      trigger: str = "scheduled") -> SyncStats:
    """
    Sync a single target: discover files, download, extract, chunk, embed, store.
    """
    label = target.label or f"{target.drive_type}:{target.drive_owner}{target.folder_path}"
    log.info(f"\n{'═'*60}")
    log.info(f"  Syncing: {label}")
    log.info(f"{'═'*60}")
    
    stats = SyncStats()
    start_time = time.time()
    run_id = None
    
    if not dry_run:
        run_id = await record_sync_run(pool, target.id, trigger)
    
    try:
        # 1. Resolve drive ID if needed
        drive_id = await graph.resolve_drive_id(target)
        if not drive_id:
            raise RuntimeError(f"Could not resolve drive for {target.drive_type}:{target.drive_owner}")
        
        if not target.drive_id and not dry_run:
            await pool.execute(
                "UPDATE chat.onedrive_sync_targets SET drive_id = $2 WHERE id = $1",
                target.id, drive_id
            )
            target.drive_id = drive_id
        
        # 2. Resolve folder ID if needed
        folder_id = target.folder_id or await graph.resolve_folder_id(drive_id, target.folder_path)
        if not folder_id:
            raise RuntimeError(f"Could not resolve folder: {target.folder_path}")
        
        if not target.folder_id and folder_id != 'root' and not dry_run:
            await pool.execute(
                "UPDATE chat.onedrive_sync_targets SET folder_id = $2 WHERE id = $1",
                target.id, folder_id
            )
        
        # 3. Get existing files for change detection
        existing = await get_existing_files(pool, target.id) if not dry_run else {}
        
        # 4. Discover files
        extension_filter = set(target.file_extensions) if target.file_extensions else SUPPORTED_EXTENSIONS
        
        use_delta = target.delta_link is not None and target.sync_mode == 'full'
        
        if use_delta:
            log.info("  Using delta sync (incremental)...")
            delta_items, new_delta_link = await graph.get_delta(
                drive_id, folder_id, target.delta_link
            )
            drive_items = []
            deleted_item_ids = set()
            
            for item in delta_items:
                if item.get("deleted"):
                    deleted_item_ids.add(item.get("id", ""))
                    continue
                if "file" not in item:
                    continue
                name = item.get("name", "")
                ext = name.rsplit('.', 1)[-1].lower() if '.' in name else ''
                if ext not in extension_filter:
                    continue
                
                created_by_info = item.get("createdBy", {}).get("user", {})
                modified_by_info = item.get("lastModifiedBy", {}).get("user", {})
                parent_ref = item.get("parentReference", {})
                parent_path = parent_ref.get("path", "")
                # Extract relative path
                if "/root:" in parent_path:
                    rel_path = parent_path.split("/root:")[-1] + "/"
                else:
                    rel_path = "/"
                
                drive_items.append(DriveItem(
                    item_id=item.get("id", ""),
                    drive_id=drive_id,
                    name=name,
                    path=f"{rel_path}{name}",
                    size=item.get("size", 0),
                    mime_type=item.get("file", {}).get("mimeType", ""),
                    extension=ext,
                    created_by=created_by_info.get("displayName"),
                    created_by_email=created_by_info.get("email"),
                    modified_by=modified_by_info.get("displayName"),
                    modified_by_email=modified_by_info.get("email"),
                    web_url=item.get("webUrl"),
                    created_at=_parse_graph_date(item.get("createdDateTime")),
                    modified_at=_parse_graph_date(item.get("lastModifiedDateTime")),
                    etag=item.get("eTag"),
                    ctag=item.get("cTag"),
                    parent_item_id=parent_ref.get("id"),
                ))
            
            # Handle deletes
            if deleted_item_ids and not dry_run:
                delete_uuids = []
                for key, existing_file in existing.items():
                    item_id = key.split(":", 1)[1]
                    if item_id in deleted_item_ids:
                        delete_uuids.append(existing_file["id"])
                if delete_uuids:
                    await mark_files_deleted(pool, delete_uuids)
                    stats.files_deleted = len(delete_uuids)
                    log.info(f"  🗑️  Marked {len(delete_uuids)} files as deleted")
        else:
            log.info("  Full crawl...")
            drive_items = await graph.crawl_folder(
                drive_id, folder_id, target.folder_path,
                target.recursive, extension_filter, target.max_file_size
            )
            new_delta_link = None
            
            # On full crawl, also get a delta link for next time
            if target.sync_mode == 'full' and not dry_run:
                _, new_delta_link = await graph.get_delta(drive_id, folder_id)
        
        stats.files_discovered = len(drive_items)
        log.info(f"  📁 Discovered {stats.files_discovered} files")
        
        if dry_run:
            for item in drive_items[:20]:
                log.info(f"    {item.path} ({item.size/1024:.0f}KB, {item.extension})")
            if len(drive_items) > 20:
                log.info(f"    ... and {len(drive_items) - 20} more")
            return stats
        
        # 5. Process each file
        for i, item in enumerate(drive_items, 1):
            key = f"{item.drive_id}:{item.item_id}"
            existing_file = existing.get(key)
            
            # Check if content changed (ctag comparison)
            if existing_file and existing_file.get("ctag") == item.ctag:
                if existing_file["status"] == "indexed":
                    stats.files_skipped += 1
                    continue
            
            is_new = existing_file is None
            
            try:
                # Upsert file record
                file_id = await upsert_file_record(pool, target, item, "downloading")
                
                # Download
                await pool.execute(
                    "UPDATE chat.onedrive_files SET status = 'downloading' WHERE id = $1",
                    file_id
                )
                content = await graph.download_file(drive_id, item.item_id)
                if not content:
                    await mark_file_failed(pool, file_id, "Download failed")
                    stats.files_failed += 1
                    continue
                
                # Extract text
                await pool.execute(
                    "UPDATE chat.onedrive_files SET status = 'extracting' WHERE id = $1",
                    file_id
                )
                text, method = extract_text(content, item.extension, item.name)
                
                if not text or not text.strip():
                    await mark_file_skipped(pool, file_id, f"No extractable text ({method})")
                    stats.files_skipped += 1
                    log.info(f"  [{i}/{len(drive_items)}] ⊘ {item.name} (no text)")
                    continue
                
                # Chunk
                await pool.execute(
                    "UPDATE chat.onedrive_files SET status = 'chunking' WHERE id = $1",
                    file_id
                )
                chunk_texts = chunker.split(text)
                if not chunk_texts:
                    await mark_file_skipped(pool, file_id, "Chunking produced no output")
                    stats.files_skipped += 1
                    continue
                
                # Contextual Retrieval — enrich chunks before embedding
                ext = item.extension or ""
                cr_context = generate_context_sync(
                    title=item.name,
                    domain="onedrive",
                    document_type=ext.lstrip(".") or "file",
                    content_preview=text[:3000],
                    case_number=None,
                )
                enriched_texts = enrich_chunks(cr_context, chunk_texts)

                # Embed in batches (using enriched texts)
                await pool.execute(
                    "UPDATE chat.onedrive_files SET status = 'embedding' WHERE id = $1",
                    file_id
                )
                all_embeddings = []
                for batch_start in range(0, len(enriched_texts), EMBEDDING_BATCH_SIZE):
                    batch = enriched_texts[batch_start:batch_start + EMBEDDING_BATCH_SIZE]
                    batch_embeds = await embedder.embed_batch(batch)
                    all_embeddings.extend(batch_embeds)
                    # Brief throttle between batches
                    if batch_start + EMBEDDING_BATCH_SIZE < len(enriched_texts):
                        await asyncio.sleep(0.5)
                
                # Write to DB
                await write_file_content(pool, file_id, text, method,
                                         chunk_texts, all_embeddings)
                
                if is_new:
                    stats.files_new += 1
                else:
                    stats.files_updated += 1
                stats.chunks_created += len(chunk_texts)
                
                action = "✚" if is_new else "↻"
                log.info(
                    f"  [{i}/{len(drive_items)}] {action} {item.name} "
                    f"→ {len(chunk_texts)} chunks ({method})"
                )
                
            except Exception as e:
                stats.files_failed += 1
                log.error(f"  [{i}/{len(drive_items)}] ✗ {item.name}: {e}")
                try:
                    if 'file_id' in dir():
                        await mark_file_failed(pool, file_id, str(e))
                except Exception:
                    pass
            
            # Progress report every 25 files
            if i % 25 == 0:
                elapsed = time.time() - start_time
                rate = i / elapsed if elapsed > 0 else 0
                eta = (len(drive_items) - i) / rate if rate > 0 else 0
                log.info(
                    f"  ── Progress: {stats.files_new} new / {stats.files_updated} updated / "
                    f"{stats.files_skipped} skipped / {stats.files_failed} failed  "
                    f"[{rate:.1f}/s, ETA {eta:.0f}s]"
                )
        
        # 6. Save delta link for next sync
        if new_delta_link and not dry_run:
            await pool.execute("""
                UPDATE chat.onedrive_sync_targets SET
                    delta_link = $2, last_sync_at = now(),
                    last_sync_files = $3, last_sync_error = NULL
                WHERE id = $1
            """, target.id, new_delta_link, stats.files_new + stats.files_updated)
        elif not dry_run:
            await pool.execute("""
                UPDATE chat.onedrive_sync_targets SET
                    last_sync_at = now(),
                    last_sync_files = $2, last_sync_error = NULL
                WHERE id = $1
            """, target.id, stats.files_new + stats.files_updated)
        
        # Update aggregate stats
        if not dry_run:
            await update_sync_target_stats(pool, target.id)
        
    except Exception as e:
        log.error(f"  ✗ Sync failed: {e}")
        if not dry_run:
            await pool.execute("""
                UPDATE chat.onedrive_sync_targets SET
                    last_sync_error = $2, last_sync_at = now()
                WHERE id = $1
            """, target.id, str(e)[:500])
        if run_id:
            await complete_sync_run(pool, run_id, stats, time.time() - start_time,
                                     error=str(e))
        raise
    
    duration = time.time() - start_time
    stats.embedding_tokens = embedder.total_tokens
    
    if run_id:
        await complete_sync_run(pool, run_id, stats, duration,
                                 delta_link=new_delta_link if not dry_run else None)
    
    return stats


# ============================================================
# Main entry points
# ============================================================

async def main(target_id: Optional[int] = None, dry_run: bool = False,
               trigger: str = "manual"):
    """Run sync for all active targets (or a specific one)."""
    pool = await asyncpg.create_pool(CHAT_DATABASE_URL, min_size=2, max_size=5)
    graph = GraphClient()
    embedder = EmbeddingClient()
    chunker = TextChunker()
    
    if not await graph.authenticate():
        log.error("Cannot authenticate to Graph API")
        return
    
    targets = await load_sync_targets(pool, target_id)
    if not targets:
        log.info("No active sync targets configured. Use --add-target to add one.")
        return
    
    log.info(f"\n{'╔'+'═'*58+'╗'}")
    log.info(f"{'║'}  OneDrive Sync — {len(targets)} target(s){' '*(39-len(str(len(targets))))}{'║'}")
    log.info(f"{'╚'+'═'*58+'╝'}")
    
    total_stats = SyncStats()
    start_time = time.time()
    
    for target in targets:
        try:
            stats = await sync_target(pool, graph, embedder, chunker,
                                       target, dry_run, trigger)
            total_stats.files_discovered += stats.files_discovered
            total_stats.files_new += stats.files_new
            total_stats.files_updated += stats.files_updated
            total_stats.files_deleted += stats.files_deleted
            total_stats.files_skipped += stats.files_skipped
            total_stats.files_failed += stats.files_failed
            total_stats.chunks_created += stats.chunks_created
        except Exception as e:
            log.error(f"Target sync failed: {e}")
    
    total_stats.embedding_tokens = embedder.total_tokens
    elapsed = time.time() - start_time
    
    log.info(f"\n{'╔'+'═'*58+'╗'}")
    log.info(f"{'║'}  ONEDRIVE SYNC COMPLETE{' '*35}{'║'}")
    log.info(f"{'╠'+'═'*58+'╣'}")
    log.info(f"{'║'}  Discovered:     {total_stats.files_discovered:<40}{'║'}")
    log.info(f"{'║'}  New:            {total_stats.files_new:<40}{'║'}")
    log.info(f"{'║'}  Updated:        {total_stats.files_updated:<40}{'║'}")
    log.info(f"{'║'}  Deleted:        {total_stats.files_deleted:<40}{'║'}")
    log.info(f"{'║'}  Skipped:        {total_stats.files_skipped:<40}{'║'}")
    log.info(f"{'║'}  Failed:         {total_stats.files_failed:<40}{'║'}")
    log.info(f"{'║'}  Chunks created: {total_stats.chunks_created:<40}{'║'}")
    log.info(f"{'║'}  Time:           {elapsed:.1f}s{' '*(38-len(f'{elapsed:.1f}s'))}{'║'}")
    log.info(f"{'║'}  Embed tokens:   {total_stats.embedding_tokens:<40}{'║'}")
    cost = total_stats.embedding_tokens * 0.00013 / 1000
    log.info(f"{'║'}  Est. cost:      ${cost:.4f}{' '*(37-len(f'${cost:.4f}'))}{'║'}")
    log.info(f"{'╚'+'═'*58+'╝'}")
    
    await embedder.close()
    await graph.close()
    await pool.close()


async def ensure_webhook_subscriptions(pool, graph):
    """Create or renew Graph change notification subscriptions for all active targets."""
    targets = await load_sync_targets(pool)
    for target in targets:
        if not target.drive_id:
            log.info(f"  Skipping webhook for target {target.id} (no drive_id yet)")
            continue
        resource = f"/drives/{target.drive_id}/root"

        # Check if we already have a valid subscription
        existing = await pool.fetchrow("""
            SELECT id, subscription_id, expiration_at
            FROM chat.onedrive_webhooks
            WHERE sync_target_id = $1 AND is_active = true
              AND expiration_at > now() + interval '1 hour'
        """, target.id)

        if existing:
            log.info(f"  Webhook for target {target.id} valid until {existing['expiration_at']}")
            # Renew if expiring within 12 hours
            if existing['expiration_at'] < datetime.now(timezone.utc) + timedelta(hours=12):
                new_expiry = datetime.now(timezone.utc) + timedelta(hours=WEBHOOK_EXPIRY_HOURS)
                ok = await graph.renew_subscription(existing['subscription_id'], new_expiry)
                if ok:
                    await pool.execute("""
                        UPDATE chat.onedrive_webhooks SET expiration_at = $2, renewed_at = now() WHERE id = $1
                    """, existing['id'], new_expiry)
                    log.info(f"  Renewed webhook for target {target.id} until {new_expiry}")
            continue

        # Create new subscription
        expiry = datetime.now(timezone.utc) + timedelta(hours=WEBHOOK_EXPIRY_HOURS)
        log.info(f"  Creating webhook for target {target.id}: {resource}")
        result = await graph.create_subscription(resource, WEBHOOK_NOTIFICATION_URL, expiry)
        if result:
            await pool.execute("""
                INSERT INTO chat.onedrive_webhooks
                    (sync_target_id, subscription_id, resource, expiration_at)
                VALUES ($1, $2, $3, $4)
                ON CONFLICT (subscription_id) DO UPDATE SET
                    expiration_at = EXCLUDED.expiration_at, is_active = true
            """, target.id, result['id'], resource, expiry)
            log.info(f"  ✓ Webhook created: {result['id']} expires {expiry}")
        else:
            log.warning(f"  ✗ Failed to create webhook for target {target.id}")


async def daemon_loop(target_id: Optional[int] = None):
    """Run sync in a continuous loop with configurable interval."""
    log.info(f"Starting OneDrive sync daemon (interval: {SYNC_INTERVAL_SECONDS}s)...")
    
    stop_event = asyncio.Event()
    webhook_setup_done = False
    
    def handle_signal(*_):
        log.info("Shutdown signal received...")
        stop_event.set()
    
    for sig in (signal.SIGTERM, signal.SIGINT):
        asyncio.get_event_loop().add_signal_handler(sig, handle_signal)
    
    while not stop_event.is_set():
        try:
            await main(target_id=target_id, trigger="scheduled")
            
            # After first successful sync, set up webhooks
            if not webhook_setup_done:
                try:
                    log.info("\n  Setting up Graph change notification webhooks...")
                    _pool = await asyncpg.create_pool(CHAT_DATABASE_URL, min_size=1, max_size=2)
                    _graph = GraphClient()
                    if await _graph.authenticate():
                        await ensure_webhook_subscriptions(_pool, _graph)
                        webhook_setup_done = True
                    await _graph.close()
                    await _pool.close()
                except Exception as e:
                    log.error(f"Webhook setup failed: {e}")
        except Exception as e:
            log.error(f"Sync cycle failed: {e}")
        
        # Wait for interval or shutdown signal
        try:
            await asyncio.wait_for(stop_event.wait(), timeout=SYNC_INTERVAL_SECONDS)
        except asyncio.TimeoutError:
            continue
    
    log.info("Daemon stopped.")


async def cli_add_target(drive_type: str, drive_owner: str,
                         folder_path: str = "/", label: str = None):
    """CLI: Add a new sync target."""
    pool = await asyncpg.create_pool(CHAT_DATABASE_URL, min_size=1, max_size=2)
    target_id = await add_sync_target(pool, drive_type, drive_owner, folder_path, label)
    log.info(f"✓ Added sync target #{target_id}: {drive_type}:{drive_owner}{folder_path}")
    if label:
        log.info(f"  Label: {label}")
    log.info(f"  Run sync with: python onedrive_sync.py --target {target_id}")
    await pool.close()


async def cli_list_targets():
    """CLI: List all sync targets."""
    pool = await asyncpg.create_pool(CHAT_DATABASE_URL, min_size=1, max_size=2)
    rows = await pool.fetch("""
        SELECT id, drive_type, drive_owner, folder_path, is_active, sync_mode,
               label, total_files, total_chunks,
               pg_size_pretty(total_bytes) as total_size,
               last_sync_at, last_sync_error
        FROM chat.onedrive_sync_targets
        ORDER BY id
    """)
    
    if not rows:
        log.info("No sync targets configured.")
        log.info("Add one with: python onedrive_sync.py --add-target user <email> /")
    else:
        log.info(f"\n{'─'*80}")
        log.info(f"  {'ID':>4}  {'Status':8}  {'Type':6}  {'Owner':<35}  {'Path':<15}  {'Files':>6}")
        log.info(f"{'─'*80}")
        for r in rows:
            status = "active" if r["is_active"] else "paused"
            log.info(
                f"  {r['id']:>4}  {status:8}  {r['drive_type']:6}  "
                f"{r['drive_owner']:<35}  {r['folder_path']:<15}  {r['total_files']:>6}"
            )
            if r["label"]:
                log.info(f"        Label: {r['label']}")
            if r["last_sync_at"]:
                log.info(f"        Last sync: {r['last_sync_at']}")
            if r["last_sync_error"]:
                log.info(f"        Error: {r['last_sync_error'][:60]}")
        log.info(f"{'─'*80}")
    
    await pool.close()


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="ACP OneDrive Sync Service")
    parser.add_argument("--target", "-t", type=int, help="Sync a specific target ID")
    parser.add_argument("--dry-run", action="store_true", help="Show what would be synced")
    parser.add_argument("--daemon", action="store_true", help="Run as continuous sync daemon")
    parser.add_argument("--interval", type=int, default=SYNC_INTERVAL_SECONDS,
                        help=f"Daemon sync interval in seconds (default: {SYNC_INTERVAL_SECONDS})")
    parser.add_argument("--list-targets", action="store_true", help="List configured sync targets")
    parser.add_argument("--add-target", nargs=3, metavar=("TYPE", "OWNER", "PATH"),
                        help="Add sync target: TYPE (user|site|group) OWNER PATH")
    parser.add_argument("--label", help="Friendly label for --add-target")
    
    args = parser.parse_args()
    
    if args.interval != SYNC_INTERVAL_SECONDS:
        SYNC_INTERVAL_SECONDS = args.interval
    
    if args.list_targets:
        asyncio.run(cli_list_targets())
    elif args.add_target:
        drive_type, drive_owner, folder_path = args.add_target
        if drive_type not in ('user', 'site', 'group'):
            print("ERROR: drive_type must be user, site, or group")
            sys.exit(1)
        asyncio.run(cli_add_target(drive_type, drive_owner, folder_path, args.label))
    elif args.daemon:
        asyncio.run(daemon_loop(target_id=args.target))
    else:
        asyncio.run(main(target_id=args.target, dry_run=args.dry_run))
