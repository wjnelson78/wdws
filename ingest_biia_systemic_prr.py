#!/usr/bin/env python3
"""
BIIA Systemic PRR Ingestion Script
===================================
Ingests ALL BIIA accommodation records (2015-2025) obtained via Public Records Act
request into the Athena Investigator Database (osint-db).

Steps:
  0. Pre-flight: upgrade embeddings from 1536 → 3072 dimensions
  1. Create investigation INV-2026-003
  2. Walk source directory, extract + classify documents
  3. Insert documents, chunks, embeddings
  4. Create timeline events for accommodation decisions
  5. Run validation queries
  6. Write summary note + ingestion report

Usage:
    python3 ingest_biia_systemic_prr.py [--dry-run] [--resume] [--skip-embed-upgrade]

    --dry-run            Extract and classify without writing to DB
    --resume             Skip files whose file_hash_sha256 already exists
    --skip-embed-upgrade Skip the embedding dimension upgrade step
"""

import argparse
import asyncio
import email as email_lib
from email import policy as email_policy
import hashlib
import json
import logging
import mimetypes
import os
import re
import subprocess
import sys
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import asyncpg
import httpx
import pytesseract
from PIL import Image

# ──────────────────────────────────────────────────────────────────────────────
# Configuration
# ──────────────────────────────────────────────────────────────────────────────

DB_URL = os.getenv(
    "OSINT_DB_URL",
    "postgresql://root:NEL2233obns@127.0.0.1:5432/osint-db"
)

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

EMBEDDING_MODEL = "text-embedding-3-large"
EMBEDDING_DIMS  = 3072
EMBEDDING_BATCH = 50

SOURCE_ROOT = Path(
    "/opt/wdws/data/dropbox/biia/"
    "09.17.2025 Nelson Request Accommodation Records 2015-2025/"
    "09.17.2025 Nelson Request Accommodation Records 2015-2025"
)

REPORT_PATH = Path("/opt/wdws/data/dropbox/biia/INGESTION_REPORT_SYSTEMIC_PRR.md")

INVESTIGATION_CASE_NUMBER = "INV-2026-003"

# Chunking
CHUNK_TARGET_CHARS  = 2500   # ≈600 tokens at ~4 chars/token
CHUNK_OVERLAP_CHARS = 200

OCR_DPI = 200

# Files / patterns to skip
SKIP_NAMES      = {".ds_store", "thumbs.db", "desktop.ini", ".ds_store"}
SKIP_PREFIXES   = ("._", "~$", "__macosx")
SKIP_EXTENSIONS = {".tmp", ".~tmp"}     # temp files
SKIP_CONTENT_EXTS = {".mp3", ".mp4", ".mov", ".avi", ".db", ".ost", ".pst", ".ini"}

# Installment folder → send date (used as fallback event_date)
INSTALLMENT_DATES: Dict[str, str] = {
    "1st": "2025-10-15",
    "2nd": "2025-11-05",
    "3rd": "2025-11-25",
    "4th": "2025-12-17",
    "5th": "2026-01-07",
    "6th": "2026-01-28",
    "7th": "2026-02-18",
}

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("biia_systemic")

if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY environment variable is required")

# Suppress noisy third-party loggers
for _noisy in ("extract_msg", "extract_msg.message", "extract_msg.attachments",
               "extract_msg.utils", "root"):
    logging.getLogger(_noisy).setLevel(logging.WARNING)


# ──────────────────────────────────────────────────────────────────────────────
# Data model
# ──────────────────────────────────────────────────────────────────────────────

@dataclass
class SystDoc:
    doc_id:             str
    file_path:          str
    filename:           str
    title:              str
    document_type:      str
    file_hash_sha256:   str
    file_size_bytes:    int
    mime_type:          str
    description:        str
    content_text:       str
    ocr_text:           Optional[str]
    tags:               List[str]
    extraction_method:  str = "unknown"
    is_duplicate:       bool = False
    existing_doc_id:    Optional[str] = None
    doc_date:           Optional[str] = None


# ──────────────────────────────────────────────────────────────────────────────
# Installment detection
# ──────────────────────────────────────────────────────────────────────────────

INSTALLMENT_PREFIXES = {
    "1st": 1, "2nd": 2, "3rd": 3, "4th": 4,
    "5th": 5, "6th": 6, "7th": 7,
}

def detect_installment(file_path: Path) -> Tuple[Optional[int], Optional[str]]:
    """Return (installment_number, send_date) from parent directory path."""
    for parent in file_path.parents:
        name_lower = parent.name.lower()
        for prefix, num in INSTALLMENT_PREFIXES.items():
            if name_lower.startswith(prefix):
                return num, INSTALLMENT_DATES.get(prefix)
    return None, None


# ──────────────────────────────────────────────────────────────────────────────
# File-skip logic
# ──────────────────────────────────────────────────────────────────────────────

def should_skip(path: Path) -> bool:
    """Return True if this file should be skipped entirely."""
    name_lower = path.name.lower()
    # macOS metadata files and temp files
    if any(name_lower.startswith(p) for p in SKIP_PREFIXES):
        return True
    if name_lower in SKIP_NAMES:
        return True
    ext = path.suffix.lower()
    if ext in SKIP_EXTENSIONS:
        return True
    if ext in SKIP_CONTENT_EXTS:
        log.debug(f"  Skipping non-text file: {path.name}")
        return True
    # __MACOSX directories
    for part in path.parts:
        if part == "__MACOSX":
            return True
    return False


# ──────────────────────────────────────────────────────────────────────────────
# Text extraction
# ──────────────────────────────────────────────────────────────────────────────

def _is_garbage(text: str) -> bool:
    """Heuristic: if >40% non-ASCII, treat as garbled."""
    if not text:
        return True
    sample = text[:2000]
    non_ascii = sum(1 for c in sample if ord(c) > 127)
    return (non_ascii / max(len(sample), 1)) > 0.40


def extract_pdf(path: Path) -> Tuple[str, str, Optional[str]]:
    """Returns (content_text, method, ocr_text_or_none)."""
    # 1. pdftotext
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", "-enc", "UTF-8", str(path), "-"],
            capture_output=True, text=True, timeout=60
        )
        text = result.stdout.strip()
        if text and len(text) > 50 and not _is_garbage(text):
            return text, "pdftotext", None
    except Exception as e:
        log.debug(f"pdftotext failed for {path.name}: {e}")

    # 2. PyMuPDF text layer
    try:
        import fitz
        doc = fitz.open(str(path))
        parts = []
        for page_num, page in enumerate(doc, 1):
            t = page.get_text("text")
            if t and t.strip():
                parts.append(f"[Page {page_num}]\n{t}")
        doc.close()
        text = "\n\n".join(parts).strip()
        if text and len(text) > 50 and not _is_garbage(text):
            return text, "pymupdf", None
    except Exception as e:
        log.debug(f"PyMuPDF text layer failed for {path.name}: {e}")

    # 3. OCR via PyMuPDF render + tesseract
    try:
        import fitz
        doc = fitz.open(str(path))
        parts = []
        for page_num, page in enumerate(doc, 1):
            mat = fitz.Matrix(OCR_DPI / 72, OCR_DPI / 72)
            pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            t = pytesseract.image_to_string(img)
            if t and t.strip():
                parts.append(f"[Page {page_num}]\n{t}")
        doc.close()
        ocr_text = "\n\n".join(parts).strip()
        return ocr_text, "tesseract_ocr", ocr_text if ocr_text else None
    except Exception as e:
        log.warning(f"OCR failed for {path.name}: {e}")

    return "", "failed", None


def extract_docx(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        import docx as docx_lib
        doc = docx_lib.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return "\n".join(paragraphs), "docx_parser", None
    except Exception as e:
        log.warning(f"DOCX extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_doc(path: Path) -> Tuple[str, str, Optional[str]]:
    """Old .doc format — try antiword, fall back to strings."""
    try:
        result = subprocess.run(
            ["antiword", str(path)],
            capture_output=True, text=True, timeout=30
        )
        if result.stdout.strip():
            return result.stdout.strip(), "antiword", None
    except Exception:
        pass
    # Fall back to docx parser (sometimes works for .doc)
    try:
        return extract_docx(path)
    except Exception:
        pass
    # Last resort: extract readable strings
    try:
        result = subprocess.run(
            ["strings", str(path)],
            capture_output=True, text=True, timeout=15
        )
        text = result.stdout.strip()
        if text and len(text) > 50:
            return text, "strings_fallback", None
    except Exception:
        pass
    return "", "failed", None


def extract_msg(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        import extract_msg as em_lib
        msg = em_lib.Message(str(path))
        parts = []
        if msg.subject:
            parts.append(f"Subject: {msg.subject}")
        if msg.sender:
            parts.append(f"From: {msg.sender}")
        if msg.to:
            parts.append(f"To: {msg.to}")
        if msg.cc:
            parts.append(f"CC: {msg.cc}")
        if msg.date:
            parts.append(f"Date: {msg.date}")
        parts.append("")
        body = msg.body or ""
        html_body = getattr(msg, "htmlBody", b"") or b""
        if body:
            parts.append(body)
        elif html_body:
            # Strip HTML tags from htmlBody
            clean = re.sub(r'<[^>]+>', ' ', html_body.decode("utf-8", errors="ignore"))
            clean = re.sub(r'\s+', ' ', clean).strip()
            parts.append(clean)
        return "\n".join(parts), "extract_msg", None
    except Exception as e:
        log.warning(f"MSG extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_eml(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        with open(path, "rb") as f:
            msg = email_lib.message_from_binary_file(f, policy=email_policy.default)
        parts = []
        for h in ("Subject", "From", "To", "CC", "Date"):
            val = msg.get(h, "")
            if val:
                parts.append(f"{h}: {val}")
        parts.append("")
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                cd = str(part.get("Content-Disposition", ""))
                if "attachment" in cd.lower():
                    continue
                if ct == "text/plain":
                    try:
                        body += part.get_content()
                    except Exception:
                        raw = part.get_payload(decode=True)
                        if raw:
                            body += raw.decode("utf-8", errors="ignore")
        else:
            try:
                body = msg.get_content()
            except Exception:
                raw = msg.get_payload(decode=True)
                if raw:
                    body = raw.decode("utf-8", errors="ignore")
        parts.append(body)
        return "\n".join(parts), "python_email", None
    except Exception as e:
        log.warning(f"EML extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_image(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        img = Image.open(str(path))
        text = pytesseract.image_to_string(img).strip()
        return text, "tesseract_ocr", text if text else None
    except Exception as e:
        log.warning(f"Image OCR failed for {path.name}: {e}")
        return "", "failed", None


def extract_xlsx(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join(str(v) for v in row if v is not None)
                if line.strip():
                    rows.append(line)
        return "\n".join(rows), "openpyxl", None
    except Exception as e:
        log.warning(f"XLSX extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_xls(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        import xlrd
        wb = xlrd.open_workbook(str(path))
        rows = []
        for sheet in wb.sheets():
            rows.append(f"[Sheet: {sheet.name}]")
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                line = "\t".join(str(cell.value) for cell in row if cell.value != "")
                if line.strip():
                    rows.append(line)
        return "\n".join(rows), "xlrd", None
    except Exception as e:
        log.warning(f"XLS extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_pptx(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n\n".join(parts), "python_pptx", None
    except Exception as e:
        log.warning(f"PPTX extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_rtf(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        from striprtf.striprtf import rtf_to_text
        raw = path.read_text(encoding="utf-8", errors="ignore")
        text = rtf_to_text(raw).strip()
        return text, "striprtf", None
    except Exception as e:
        log.warning(f"RTF extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_htm(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        raw = path.read_text(encoding="utf-8", errors="ignore")
        text = re.sub(r'<[^>]+>', ' ', raw)
        text = re.sub(r'&nbsp;', ' ', text)
        text = re.sub(r'&amp;', '&', text)
        text = re.sub(r'&lt;', '<', text)
        text = re.sub(r'&gt;', '>', text)
        text = re.sub(r'\s+', ' ', text).strip()
        return text, "html_strip", None
    except Exception as e:
        log.warning(f"HTML extraction failed for {path.name}: {e}")
        return "", "failed", None


def extract_text_file(path: Path) -> Tuple[str, str, Optional[str]]:
    try:
        return path.read_text(encoding="utf-8", errors="ignore"), "text_file", None
    except Exception as e:
        log.warning(f"Text file read failed for {path.name}: {e}")
        return "", "failed", None


def extract_content(path: Path) -> Tuple[str, str, Optional[str]]:
    """Dispatch to appropriate extractor. Returns (content_text, method, ocr_text_or_none)."""
    ext = path.suffix.lower()
    if ext in (".pdf",):
        return extract_pdf(path)
    elif ext == ".docx":
        return extract_docx(path)
    elif ext in (".doc", ".docm"):
        return extract_doc(path)
    elif ext == ".msg":
        return extract_msg(path)
    elif ext in (".eml", ".emlx"):
        return extract_eml(path)
    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"):
        return extract_image(path)
    elif ext in (".xlsx",):
        return extract_xlsx(path)
    elif ext in (".xls",):
        return extract_xls(path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    elif ext == ".rtf":
        return extract_rtf(path)
    elif ext in (".html", ".htm"):
        return extract_htm(path)
    elif ext in (".txt", ".csv", ".log"):
        return extract_text_file(path)
    else:
        # Best-effort text read for unknown types
        return extract_text_file(path)


# ──────────────────────────────────────────────────────────────────────────────
# Classification
# ──────────────────────────────────────────────────────────────────────────────

JUDGE_PATTERNS = [
    (re.compile(r'Timothy\s+M\.?\s+Blood', re.I),        'Timothy M. Blood',  'judge-blood'),
    (re.compile(r'\bTim\s+Blood\b', re.I),               'Timothy M. Blood',  'judge-blood'),
    (re.compile(r'\bJudge\s+Blood\b', re.I),             'Timothy M. Blood',  'judge-blood'),
    (re.compile(r'\bIAJ\s+Blood\b', re.I),               'Timothy M. Blood',  'judge-blood'),
    (re.compile(r'Ann\s+(?:M\.?\s+)?Dodge', re.I),       'Ann Dodge',         'judge-dodge'),
    (re.compile(r'\bJudge\s+Dodge\b', re.I),             'Ann Dodge',         'judge-dodge'),
    (re.compile(r'\bIAJ\s+Dodge\b', re.I),               'Ann Dodge',         'judge-dodge'),
    (re.compile(r'Brian\s+(?:C\.?\s+)?Watkins', re.I),   'Brian Watkins',     'judge-watkins'),
    (re.compile(r'\bJudge\s+Watkins\b', re.I),           'Brian Watkins',     'judge-watkins'),
    (re.compile(r'\bIAJ\s+Watkins\b', re.I),             'Brian Watkins',     'judge-watkins'),
    (re.compile(r'Janice\s+(?:K\.?\s+)?Rosen', re.I),   'Janice Rosen',      'judge-rosen'),
    (re.compile(r'Mark\s+(?:A\.?\s+)?Janicke', re.I),   'Mark Janicke',      'judge-janicke'),
    (re.compile(r'Chief\s+Judge\s+Janicke', re.I),       'Mark Janicke',      'judge-janicke'),
]

# Generic judge extraction for judges not in the above list
GENERIC_JUDGE_RE = re.compile(
    r'\b(?:IAJ|Industrial\s+Appeals\s+Judge|Judge|Presiding\s+Judge)\s+'
    r'([A-Z][a-z]{2,20}(?:\s+[A-Z]\.?)?\s+[A-Z][a-z]{2,20})',
    re.M
)

ACCOMM_TYPE_PATTERNS = {
    "email_filing":         re.compile(r'email\s+fil(ing|e)|electronic\s+fil(ing|e)|e-?fil(ing|e)', re.I),
    "email_service":        re.compile(r'email\s+service|service\s+by\s+email|electronic\s+service', re.I),
    "remote_hearing":       re.compile(r'remote\s+hear|telephon(ic|e)\s+hear|video\s+(hear|conference)|telephone\s+conference', re.I),
    "extended_time":        re.compile(r'extend(ed)?\s+time|additional\s+time|time\s+extension|extra\s+time', re.I),
    "interpreter":          re.compile(r'interpreter|ASL|sign\s+language|translation|language\s+access|foreign\s+language', re.I),
    "physical_access":      re.compile(r'physical\s+access|wheelchair|handicap|accessible\s+park|ramp|elevator|ADA\s+access', re.I),
    "counsel_appointment":  re.compile(r'appoint(ment)?\s+of\s+counsel|appointed\s+counsel|pro\s+se\s+counsel|legal\s+representation', re.I),
    "assistive_technology": re.compile(r'assistive\s+tech|screen\s+reader|magnif(ication|ier)|text.to.speech|voice\s+recogn', re.I),
    "large_print":          re.compile(r'large\s+print|enlarged\s+print|large\s+font|font\s+size', re.I),
    "modified_procedure":   re.compile(r'modified\s+(hearing|procedure|format)|adjust(ed)?\s+format', re.I),
    "stay_of_proceedings":  re.compile(r'stay\s+of\s+proceed|stay\s+the\s+case|postpone|continuance', re.I),
}

ACCOMM_OUTCOME_PATTERNS = {
    "granted":    re.compile(r'\b(grant(ed|ing|s)?|allow(ed|ing)?|approv(ed|ing)?)\b.{0,120}(accommodation|request)', re.I),
    "denied":     re.compile(r'\b(den(ied|y|ying|ies)|declin(ed|ing)|refus(ed|ing)|reject(ed|ing))\b.{0,120}(accommodation|request)', re.I),
    "partial":    re.compile(r'partial(ly)?\s+(grant(ed)?|allow(ed)?|approv(ed)?)', re.I),
}

# Disability type detection patterns
DISABILITY_PATTERNS = {
    "disability-cognitive-neurological": re.compile(
        r'traumatic\s+brain|TBI|brain\s+injur|long\s+covid|post.covid|neurocog|cognitive\s+impairment'
        r'|neurolog|brain\s+fog|processing\s+speed|executive\s+function|memory\s+impairment'
        r'|concussion|dementia|alzheimer|ADHD|attention\s+deficit|autism|ASD|asperger'
        r'|intellectual\s+disabilit',
        re.I
    ),
    "disability-physical": re.compile(
        r'wheelchair|mobility\s+impairment|physical\s+disabilit|orthopedic|spinal|back\s+injur'
        r'|amputat|prosthetic|walking|crutch(es)?|cane|walker|musculoskeletal|chronic\s+pain'
        r'|fibromyalgia|arthritis|joint|limited\s+mobility',
        re.I
    ),
    "disability-sensory-visual": re.compile(
        r'\bblind(ness|ed)?\b|low\s+vision|visual\s+impairment|sight\s+impairment|Braille'
        r'|macular|glaucoma|retinal',
        re.I
    ),
    "disability-sensory-hearing": re.compile(
        r'\bdeaf(ness)?\b|hard\s+of\s+hearing|hearing\s+impairment|hearing\s+loss|ASL'
        r'|sign\s+language|cochlear',
        re.I
    ),
    "disability-mental-health": re.compile(
        r'PTSD|post.traumatic|anxiety\s+disorder|panic\s+disorder|depression|depressive'
        r'|bipolar|schizophreni|psychiat|mental\s+health|mental\s+illness|OCD|phobia'
        r'|agoraphobia|generalized\s+anxiety',
        re.I
    ),
    "disability-developmental": re.compile(
        r'developmental\s+disabilit|Down\s+syndrome|cerebral\s+palsy|intellectual\s+disability'
        r'|birth\s+defect|pervasive\s+developmental',
        re.I
    ),
}

# Content-based document type detection (highest priority)
CONTENT_DOCTYPE_PATTERNS = [
    (re.compile(r'ORDER\s+(ON|GRANTING|DENYING|RE:|REGARDING|ADDRESSING)\s+(?:REQUEST\s+FOR\s+)?ACCOMMODATION', re.I), "accommodation_order"),
    (re.compile(r'ORDER\s+ON\s+.{0,40}ACCOMMODATION', re.I), "accommodation_order"),
    (re.compile(r'ORDER\s+STAYING', re.I), "accommodation_order"),
    (re.compile(r'GENERAL ORDER\s+NO\.?\s*\d', re.I), "accommodation_order"),
    (re.compile(r'(?:BIIA|Board).{0,40}ACCOMMODATION\s+REQUEST', re.I), "accommodation_request"),
    (re.compile(r'I\s+(?:hereby\s+)?request\s+(an?\s+)?accommodation', re.I), "accommodation_request"),
    (re.compile(r'ADA\s+ACCOMMODATION\s+REQUEST\s+FORM', re.I), "accommodation_form"),
    (re.compile(r'REQUEST\s+FOR\s+REASONABLE\s+ACCOMMODATION', re.I), "accommodation_request"),
    (re.compile(r'ACCOMMODATION\s+REQUEST\s+FORM', re.I), "accommodation_form"),
    (re.compile(r'(?:BIIA|Board).{0,60}(?:POLICY|PROCEDURE|GUIDELINES?|PROTOCOL|SOP)\s+(?:FOR|ON|RE|REGARDING)\s+(?:ADA|ACCOMMODATION)', re.I), "accommodation_policy"),
    (re.compile(r'STANDARD\s+OPERATING\s+PROCEDURE.{0,60}ACCOMMODATION', re.I), "accommodation_policy"),
    (re.compile(r'ACCOMMODATION.*TRAINING', re.I), "accommodation_policy"),
    (re.compile(r'CERTIFIED\s+MAIL(?:ING)?', re.I), "mailroom_record"),
    (re.compile(r'USPS\s+(?:Return\s+Receipt|Tracking)', re.I), "mailroom_record"),
    (re.compile(r'BAMS\s+Mailing', re.I), "mailroom_record"),
    (re.compile(r'voicemail', re.I), "phone_record"),
    (re.compile(r'EXEMPTION\s+KEY|REDACTION\s+KEY|EXEMPTION\s+LOG', re.I), "exemption_key"),
    (re.compile(r'Diagnosis|Medical\s+Record|Clinical\s+Note|Physician|Doctor', re.I), "medical_documentation"),
]

# Folder-name hints
FOLDER_DOCTYPE = {
    "order":                    "accommodation_order",
    "accommodation order":      "accommodation_order",
    "accommodation request":    "accommodation_request",
    "accommodation request form": "accommodation_form",
    "forms":                    "accommodation_form",
    "e-mailed correspondence":  "external_communication",
    "email correspondence":     "external_communication",
    "mailed correspondence":    "external_communication",
    "mailroom":                 "mailroom_record",
    "mailroom records":         "mailroom_record",
    "claim information":        "other",
    "policies":                 "accommodation_policy",
    "policy":                   "accommodation_policy",
    "procedures":               "accommodation_policy",
    "training":                 "accommodation_policy",
    "subcommittee":             "accommodation_policy",
    "old forms":                "accommodation_form",
    "voice mail":               "phone_record",
    "voicemail":                "phone_record",
    "medical":                  "medical_documentation",
}

# Filename hints
FILE_BASED_DOCTYPE = {
    "order":                "accommodation_order",
    "ord_":                 "accommodation_order",
    "template":             "accommodation_policy",
    "policy":               "accommodation_policy",
    "form":                 "accommodation_form",
    "request":              "accommodation_request",
    "case log":             "internal_communication",
    "active cases":         "internal_communication",
    "distribution list":    "internal_communication",
    "dist list":            "internal_communication",
    "letter":               "external_communication",
    "certified mail":       "mailroom_record",
    "usps":                 "mailroom_record",
    "return receipt":       "mailroom_record",
    "voicemail":            "phone_record",
    "phone":                "phone_record",
    "medical":              "medical_documentation",
    "diagnosis":            "medical_documentation",
    "exemption":            "exemption_key",
}


def classify_document(path: Path, text: str) -> str:
    """Classify document type using content, filename, and folder signals."""
    filename_lower = path.name.lower()
    text_sample = text[:3000] if text else ""

    # Content-based (highest priority)
    for pattern, dtype in CONTENT_DOCTYPE_PATTERNS:
        if pattern.search(text_sample):
            return dtype

    # Folder path hints
    for parent in path.parents:
        folder_lower = parent.name.lower()
        for folder_key, dtype in FOLDER_DOCTYPE.items():
            if folder_key in folder_lower:
                return dtype

    # Filename hints
    for name_key, dtype in FILE_BASED_DOCTYPE.items():
        if name_key in filename_lower:
            return dtype

    # Extension-based fallbacks
    ext = path.suffix.lower()
    if ext == ".msg":
        return "internal_communication"
    if ext in (".eml", ".emlx"):
        return "internal_communication"

    return "other"


# ──────────────────────────────────────────────────────────────────────────────
# Tag building
# ──────────────────────────────────────────────────────────────────────────────

def build_tags(path: Path, text: str, doc_type: str) -> Tuple[List[str], Optional[str], Optional[str]]:
    """Build tags array.
    Returns (tags, judge_name_or_none, outcome_or_none)."""
    text_sample = text[:6000] if text else ""
    tags = ["biia-prr-systemic", "accommodation", "prr-2015-2025"]

    # Document type tag
    type_tag_map = {
        "accommodation_request":    "accommodation-request",
        "accommodation_order":      "accommodation-order",
        "accommodation_policy":     "accommodation-policy",
        "accommodation_form":       "accommodation-request",
        "internal_communication":   "internal-communication",
        "external_communication":   "external-communication",
        "phone_record":             "internal-communication",
        "mailroom_record":          "internal-communication",
        "medical_documentation":    "medical-documentation",
        "exemption_key":            "accommodation-policy",
        "other":                    None,
    }
    dtype_tag = type_tag_map.get(doc_type)
    if dtype_tag:
        tags.append(dtype_tag)

    # Outcome detection
    outcome = None
    if doc_type in ("accommodation_order", "accommodation_request", "external_communication"):
        # Check partial first (most specific)
        if ACCOMM_OUTCOME_PATTERNS["partial"].search(text_sample):
            tags.append("outcome-partial")
            outcome = "partial"
        elif ACCOMM_OUTCOME_PATTERNS["granted"].search(text_sample):
            tags.append("outcome-granted")
            outcome = "granted"
        elif ACCOMM_OUTCOME_PATTERNS["denied"].search(text_sample):
            tags.append("outcome-denied")
            outcome = "denied"

    # Disability type
    disability_found = False
    for dis_tag, pattern in DISABILITY_PATTERNS.items():
        if pattern.search(text_sample):
            tags.append(dis_tag)
            disability_found = True
    if not disability_found and re.search(r'disabilit|impairment|condition', text_sample, re.I):
        tags.append("disability-unspecified")

    # Judge detection
    judge_name = None
    for pattern, name, tag in JUDGE_PATTERNS:
        if pattern.search(text_sample):
            tags.append(tag)
            judge_name = name
            break
    if not judge_name:
        # Try generic judge extraction
        m = GENERIC_JUDGE_RE.search(text_sample)
        if m:
            raw_name = m.group(1).strip()
            # Skip generic "the Board" type matches
            skip_words = {"the", "this", "an", "a", "any", "industrial", "appeals"}
            name_parts = raw_name.lower().split()
            if not any(p in skip_words for p in name_parts):
                last = raw_name.split()[-1].lower()
                tags.append(f"judge-{last}")
                judge_name = raw_name

    # Accommodation type tags
    for accomm_type, pattern in ACCOMM_TYPE_PATTERNS.items():
        if pattern.search(text_sample):
            tag_map = {
                "email_filing":         "accom-email-filing",
                "email_service":        "accom-email-filing",
                "remote_hearing":       "accom-remote-hearing",
                "extended_time":        "accom-extended-time",
                "interpreter":          "accom-interpreter",
                "physical_access":      "accom-physical-access",
                "counsel_appointment":  "accom-counsel",
                "assistive_technology": "accom-assistive-tech",
                "large_print":          "accom-extended-time",
                "modified_procedure":   "accom-remote-hearing",
                "stay_of_proceedings":  "accom-extended-time",
            }
            t = tag_map.get(accomm_type)
            if t and t not in tags:
                tags.append(t)

    # Redaction detection
    if re.search(r'\[REDACTED\]|\[WITHHELD\]|Exemption\s+\d+|RCW\s+42\.56', text_sample, re.I):
        tags.append("redacted")

    # Nelson-related detection
    nelson_names = ["william nelson", "w. nelson", "william james nelson", "nelson, william",
                    "25-18153", "sx-21824", "25 18153"]
    if any(n in text[:8000].lower() for n in nelson_names):
        tags.append("nelson-related")

    return tags, judge_name, outcome


# ──────────────────────────────────────────────────────────────────────────────
# Date extraction
# ──────────────────────────────────────────────────────────────────────────────

MONTH_MAP = {
    "january": "01", "february": "02", "march": "03", "april": "04",
    "may": "05", "june": "06", "july": "07", "august": "08",
    "september": "09", "october": "10", "november": "11", "december": "12",
}


def extract_doc_date(text: str, filename: str = "") -> Optional[str]:
    """Extract document date as YYYY-MM-DD from content or filename."""
    # ISO date from filename
    m = re.search(r'(\d{4})[-._](\d{2})[-._](\d{2})', filename)
    if m:
        y, mo, d = m.groups()
        if 2013 <= int(y) <= 2026 and 1 <= int(mo) <= 12 and 1 <= int(d) <= 31:
            return f"{y}-{mo}-{d}"

    # Filename date like "9.12.25" or "09.12.24"
    m = re.search(r'(\d{1,2})\.(\d{2})\.(\d{2})\b', filename)
    if m:
        mo, d, y_short = m.groups()
        y = int(y_short) + 2000
        if 2013 <= y <= 2030:
            return f"{y}-{mo.zfill(2)}-{d.zfill(2)}"

    sample = text[:8000] if text else ""

    # Month name full: "September 12, 2025"
    m = re.search(
        r'(January|February|March|April|May|June|July|August|September|October|November|December)'
        r'\s+(\d{1,2}),?\s+(20\d{2})',
        sample, re.I
    )
    if m:
        mo_name, d, y = m.groups()
        mo = MONTH_MAP.get(mo_name.lower(), "00")
        if mo != "00" and 2013 <= int(y) <= 2026:
            return f"{y}-{mo}-{d.zfill(2)}"

    # MM/DD/YYYY
    m = re.search(r'(\d{1,2})/(\d{1,2})/(20\d{2})', sample)
    if m:
        mo, d, y = m.groups()
        if 2013 <= int(y) <= 2026:
            return f"{y}-{mo.zfill(2)}-{d.zfill(2)}"

    # YYYY-MM-DD in content
    m = re.search(r'(20[1-2]\d)-(\d{2})-(\d{2})', sample)
    if m:
        y, mo, d = m.groups()
        if 2013 <= int(y) <= 2026:
            return f"{y}-{mo}-{d}"

    return None


def extract_msg_date(path: Path) -> Optional[str]:
    try:
        import extract_msg as em_lib
        msg = em_lib.Message(str(path))
        d = msg.date
        if d:
            m = re.search(r'(\d{4})-(\d{2})-(\d{2})', str(d))
            if m:
                return f"{m.group(1)}-{m.group(2)}-{m.group(3)}"
    except Exception:
        pass
    return None


# ──────────────────────────────────────────────────────────────────────────────
# Title + description
# ──────────────────────────────────────────────────────────────────────────────

def build_title(path: Path, text: str) -> str:
    """Derive a clean title from content or filename."""
    if text:
        top = text[:2000].strip()
        for pat in [
            re.compile(r'^(ORDER\s+[A-Z][A-Z\s/,:\-]{5,100}?)(?:\n|$)', re.M),
            re.compile(r'^(ACCOMMODATION\s+[A-Z][A-Z\s/,:\-]{5,80}?)(?:\n|$)', re.M),
            re.compile(r'^(BIIA\s+[A-Z][A-Z\s]{5,80}?)(?:\n|$)', re.M),
        ]:
            m = pat.search(top)
            if m:
                candidate = m.group(1).strip()
                if 10 < len(candidate) < 150:
                    return candidate

    # Fall back to cleaned filename
    stem = path.stem
    stem = re.sub(r'^[\d./_\- ]+', '', stem).strip()   # remove leading date/num prefixes
    stem = re.sub(r'[-_]+', ' ', stem)
    stem = re.sub(r'\s+', ' ', stem).strip()
    return stem if stem else path.name


def build_description(
    path: Path,
    text: str,
    doc_type: str,
    tags: List[str],
    judge_name: Optional[str],
    outcome: Optional[str],
) -> str:
    """Generate a 1-3 sentence description from extracted metadata."""
    parts = []

    # What type of document
    type_labels = {
        "accommodation_request":    "Accommodation request",
        "accommodation_order":      "BIIA accommodation order",
        "accommodation_policy":     "BIIA accommodation policy/procedure document",
        "accommodation_form":       "ADA accommodation request form",
        "internal_communication":   "Internal BIIA communication",
        "external_communication":   "External communication (BIIA ↔ party/attorney)",
        "phone_record":             "Phone/voicemail record",
        "mailroom_record":          "Mailroom/mailing record",
        "medical_documentation":    "Medical documentation submitted with accommodation request",
        "exemption_key":            "Universal Exemption Key (BIIA redaction codes reference document)",
        "other":                    "Document",
    }
    doc_label = type_labels.get(doc_type, "Document")

    # Judge
    judge_part = f" by {judge_name}" if judge_name else ""

    # Docket from folder name
    docket = None
    for parent in path.parents:
        m = re.search(r'Docket\s+No\.?\s*([\d\s]+)', parent.name, re.I)
        if not m:
            m = re.search(r'(\d{2}\s*\d{4,5})', parent.name)
        if m:
            docket = m.group(1).strip().replace(" ", " ")
            break

    docket_part = f" Docket {docket}." if docket else ""

    # Accommodation type
    accomm_parts = []
    if "accom-email-filing" in tags:
        accomm_parts.append("email filing")
    if "accom-remote-hearing" in tags:
        accomm_parts.append("remote/telephone hearing")
    if "accom-interpreter" in tags:
        accomm_parts.append("interpreter services")
    if "accom-extended-time" in tags:
        accomm_parts.append("extended time")
    if "accom-physical-access" in tags:
        accomm_parts.append("physical accessibility")
    if "accom-counsel" in tags:
        accomm_parts.append("appointment of counsel")
    if "accom-assistive-tech" in tags:
        accomm_parts.append("assistive technology")
    accomm_text = ", ".join(accomm_parts) if accomm_parts else "accommodation"

    # Disability type
    dis_parts = []
    if "disability-cognitive-neurological" in tags:
        dis_parts.append("cognitive/neurological")
    if "disability-physical" in tags:
        dis_parts.append("physical/mobility")
    if "disability-sensory-visual" in tags:
        dis_parts.append("visual")
    if "disability-sensory-hearing" in tags:
        dis_parts.append("deaf/hard of hearing")
    if "disability-mental-health" in tags:
        dis_parts.append("mental health")
    if "disability-developmental" in tags:
        dis_parts.append("developmental")
    if "disability-unspecified" in tags and not dis_parts:
        dis_parts.append("unspecified")
    dis_text = "/".join(dis_parts) if dis_parts else None

    # Outcome
    outcome_labels = {
        "granted": "Accommodation GRANTED.",
        "denied":  "Accommodation DENIED.",
        "partial": "Accommodation PARTIALLY granted.",
    }
    outcome_text = outcome_labels.get(outcome, "")

    # Build sentence 1
    sent1 = f"{doc_label}{judge_part}.{docket_part}"
    if accomm_text != "accommodation" or outcome_text:
        sent1 = f"{doc_label}{judge_part} — {accomm_text} request.{docket_part}"

    parts.append(sent1)

    # Build sentence 2 (disability + outcome)
    if dis_text and outcome_text:
        parts.append(f"Disability type: {dis_text}. {outcome_text}")
    elif dis_text:
        parts.append(f"Disability type: {dis_text}.")
    elif outcome_text:
        parts.append(outcome_text)

    # Nelson-related note
    if "nelson-related" in tags:
        parts.append("Relates to William Nelson (Docket 25-18153).")

    return " ".join(parts).strip()


# ──────────────────────────────────────────────────────────────────────────────
# Chunking
# ──────────────────────────────────────────────────────────────────────────────

def chunk_text(text: str) -> List[Tuple[str, int, int]]:
    """Split text into overlapping chunks.
    Returns list of (chunk_text, start_char, end_char).
    """
    if not text:
        return []
    if len(text) <= CHUNK_TARGET_CHARS:
        return [(text, 0, len(text))]

    paragraphs = re.split(r'\n{2,}', text)
    chunks: List[Tuple[str, int, int]] = []
    current = ""
    current_start = 0
    char_offset = 0

    for para in paragraphs:
        if not para.strip():
            char_offset += len(para) + 2  # +2 for the \n\n split
            continue
        candidate = (current + "\n\n" + para).strip() if current else para.strip()
        if len(candidate) <= CHUNK_TARGET_CHARS:
            if not current:
                current_start = char_offset
            current = candidate
        else:
            if current:
                end = current_start + len(current)
                chunks.append((current, current_start, end))
                # Overlap
                overlap_text = current[-CHUNK_OVERLAP_CHARS:]
                current_start = max(0, end - CHUNK_OVERLAP_CHARS)
                current = (overlap_text + "\n\n" + para.strip()).strip()
            else:
                # Para is too large — split on sentences
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sent in sentences:
                    candidate = (current + " " + sent).strip() if current else sent.strip()
                    if len(candidate) <= CHUNK_TARGET_CHARS:
                        if not current:
                            current_start = char_offset
                        current = candidate
                    else:
                        if current:
                            end = current_start + len(current)
                            chunks.append((current, current_start, end))
                        current_start = char_offset
                        current = sent.strip()
        char_offset += len(para) + 2

    if current:
        end = current_start + len(current)
        chunks.append((current, current_start, end))

    return [(c, s, e) for c, s, e in chunks if c.strip()]


def count_tokens_approx(text: str) -> int:
    return max(1, len(text) // 4)


# ──────────────────────────────────────────────────────────────────────────────
# OpenAI embedding client
# ──────────────────────────────────────────────────────────────────────────────

class EmbeddingClient:
    def __init__(self):
        self.client = httpx.AsyncClient(timeout=120.0)
        self.total_tokens = 0
        self.total_requests = 0

    async def embed(self, texts: List[str]) -> List[List[float]]:
        if not texts:
            return []
        truncated = [t[:30000] for t in texts]
        resp = await self.client.post(
            "https://api.openai.com/v1/embeddings",
            headers={"Authorization": f"Bearer {OPENAI_API_KEY}"},
            json={"input": truncated, "model": EMBEDDING_MODEL, "dimensions": EMBEDDING_DIMS},
        )
        resp.raise_for_status()
        data = resp.json()
        self.total_tokens += data.get("usage", {}).get("total_tokens", 0)
        self.total_requests += 1
        return [d["embedding"] for d in data["data"]]

    async def embed_batched(self, texts: List[str]) -> List[List[float]]:
        """Embed in batches of EMBEDDING_BATCH."""
        results = []
        for i in range(0, len(texts), EMBEDDING_BATCH):
            batch = texts[i:i + EMBEDDING_BATCH]
            batch_results = await self.embed(batch)
            results.extend(batch_results)
        return results

    async def close(self):
        await self.client.aclose()


def vec_literal(embedding: List[float]) -> str:
    return "[" + ",".join(f"{v:.8f}" for v in embedding) + "]"


# ──────────────────────────────────────────────────────────────────────────────
# Embedding dimension upgrade
# ──────────────────────────────────────────────────────────────────────────────

async def upgrade_embeddings(pool: asyncpg.Pool, embedder: EmbeddingClient, dry_run: bool) -> Dict:
    """Upgrade document_chunks and entity_embeddings from 1536 → 3072 dimensions."""
    log.info("=== STEP 0: Upgrading embeddings from 1536 → 3072 dimensions ===")
    stats = {
        "chunks_processed": 0,
        "entity_embs_processed": 0,
        "tokens_used_before": embedder.total_tokens,
        "skipped": False,
    }

    # Check if already upgraded
    row = await pool.fetchrow(
        "SELECT vector_dims(embedding) as dims FROM document_chunks WHERE embedding IS NOT NULL LIMIT 1"
    )
    if row and row["dims"] == EMBEDDING_DIMS:
        log.info("  document_chunks already at 3072 dims — skipping upgrade")
        stats["skipped"] = True
        return stats

    if dry_run:
        log.info("  [DRY RUN] Would upgrade 1536→3072 dims on document_chunks + entity_embeddings")
        return stats

    async with pool.acquire() as conn:
        async with conn.transaction():
            log.info("  Dropping HNSW indexes...")
            await conn.execute("DROP INDEX IF EXISTS idx_document_chunks_hnsw")
            await conn.execute("DROP INDEX IF EXISTS idx_entity_embeddings_hnsw")

            log.info("  Altering document_chunks.embedding to vector(3072)...")
            await conn.execute(
                "ALTER TABLE document_chunks ALTER COLUMN embedding DROP NOT NULL"
            )
            await conn.execute(
                "UPDATE document_chunks SET embedding = NULL"
            )
            await conn.execute(
                "ALTER TABLE document_chunks ALTER COLUMN embedding TYPE vector(3072)"
            )

            log.info("  Altering entity_embeddings.embedding to vector(3072)...")
            await conn.execute(
                "ALTER TABLE entity_embeddings ALTER COLUMN embedding DROP NOT NULL"
            )
            await conn.execute(
                "UPDATE entity_embeddings SET embedding = NULL"
            )
            await conn.execute(
                "ALTER TABLE entity_embeddings ALTER COLUMN embedding TYPE vector(3072)"
            )

    log.info("  Schema upgraded. Re-embedding document_chunks...")

    # Re-embed document_chunks
    chunks = await pool.fetch(
        "SELECT id, chunk_text FROM document_chunks WHERE chunk_text IS NOT NULL ORDER BY id"
    )
    log.info(f"  Found {len(chunks)} chunks to re-embed")

    chunk_texts = [r["chunk_text"] for r in chunks]
    chunk_ids   = [r["id"] for r in chunks]
    embeddings  = await embedder.embed_batched(chunk_texts)

    async with pool.acquire() as conn:
        async with conn.transaction():
            for cid, emb in zip(chunk_ids, embeddings):
                await conn.execute(
                    "UPDATE document_chunks SET embedding = $1::vector WHERE id = $2",
                    vec_literal(emb), cid
                )
    stats["chunks_processed"] = len(chunks)
    log.info(f"  ✓ Re-embedded {len(chunks)} document chunks")

    # Re-embed entity_embeddings
    entity_rows = await pool.fetch(
        "SELECT id, source_text FROM entity_embeddings WHERE source_text IS NOT NULL ORDER BY id"
    )
    log.info(f"  Found {len(entity_rows)} entity embeddings to re-embed")

    ent_texts = [r["source_text"] for r in entity_rows]
    ent_ids   = [r["id"] for r in entity_rows]
    ent_embeddings = await embedder.embed_batched(ent_texts)

    async with pool.acquire() as conn:
        async with conn.transaction():
            for eid, emb in zip(ent_ids, ent_embeddings):
                await conn.execute(
                    "UPDATE entity_embeddings SET embedding = $1::vector, model = $2 WHERE id = $3",
                    vec_literal(emb), EMBEDDING_MODEL, eid
                )
    stats["entity_embs_processed"] = len(entity_rows)
    log.info(f"  ✓ Re-embedded {len(entity_rows)} entity embeddings")

    # Recreate HNSW indexes
    log.info("  Recreating HNSW indexes...")
    async with pool.acquire() as conn:
        await conn.execute(
            "CREATE INDEX CONCURRENTLY IF NOT EXISTS idx_document_chunks_hnsw "
            "ON document_chunks USING hnsw (embedding vector_cosine_ops) "
            "WITH (m=16, ef_construction=64)"
        )
        await conn.execute(
            "CREATE INDEX CONCURRENTLY IF NOT EXISTS idx_entity_embeddings_hnsw "
            "ON entity_embeddings USING hnsw (embedding vector_cosine_ops) "
            "WITH (m=16, ef_construction=64)"
        )
    log.info("  ✓ HNSW indexes recreated")

    # Log embedding job
    tokens_used = embedder.total_tokens - stats["tokens_used_before"]
    total_processed = stats["chunks_processed"] + stats["entity_embs_processed"]
    await pool.execute("""
        INSERT INTO embedding_jobs
            (job_type, status, entities_processed, tokens_used, model, completed_at)
        VALUES ($1, $2, $3, $4, $5, NOW())
    """,
        "migration_small_to_large",
        "completed",
        total_processed,
        tokens_used,
        EMBEDDING_MODEL,
    )
    log.info(f"  ✓ Embedding upgrade complete. Tokens used: {tokens_used:,}")
    return stats


# ──────────────────────────────────────────────────────────────────────────────
# Investigation
# ──────────────────────────────────────────────────────────────────────────────

async def get_or_create_investigation(pool: asyncpg.Pool, dry_run: bool) -> str:
    """Return UUID of INV-2026-003, creating it if needed."""
    row = await pool.fetchrow(
        "SELECT id::text FROM investigations WHERE case_number = $1",
        INVESTIGATION_CASE_NUMBER
    )
    if row:
        log.info(f"  Investigation {INVESTIGATION_CASE_NUMBER} already exists: {row['id']}")
        return row["id"]

    if dry_run:
        fake_id = str(uuid.uuid4())
        log.info(f"  [DRY RUN] Would create investigation {INVESTIGATION_CASE_NUMBER}")
        return fake_id

    new_id = str(uuid.uuid4())
    await pool.execute("""
        INSERT INTO investigations (
            id, case_number, title, description, status, priority,
            lead_investigator, tags, metadata
        ) VALUES (
            $1::uuid, $2, $3, $4, $5::investigation_status, $6, $7, $8, $9::jsonb
        )
    """,
        new_id,
        INVESTIGATION_CASE_NUMBER,
        "Investigation: BIIA Systemic ADA Accommodation Practices (2015-2025)",
        (
            "Analysis of Washington State Board of Industrial Insurance Appeals (BIIA) accommodation "
            "records obtained via Public Records Act request filed September 16, 2025. Dataset covers "
            "ALL litigant accommodation requests from January 1, 2015 through September 17, 2025. "
            "Purpose: identify systemic patterns in how BIIA handles ADA/Title II accommodation "
            "requests, with focus on disparities between visible/physical disabilities and "
            "invisible/neurological disabilities. Supports Nelson v. BIIA federal litigation and "
            "potential media investigation by KING 5."
        ),
        "active",
        1,
        "wnelson",
        [
            "biia", "ada", "accommodation", "systemic-analysis", "prr",
            "disability-rights", "gr33", "title-ii", "pattern-analysis"
        ],
        json.dumps({
            "prr_request_date": "2025-09-16",
            "prr_scope": "all_litigants_2015_2025",
            "biia_searches": {
                "search_1_accommodations_folder": "2900+ files",
                "search_2_zdrive": "14500+ files",
                "search_5_webdocs": "200+ files"
            },
            "box_link": "https://biia.box.com/s/t9sdzyk2v9dyv0jc6m4nvtayt892ef2g",
            "source_folder": "09.17.2025 Nelson Request Accommodation Records 2015-2025",
            "withheld_records": "103 crime victim records withheld",
            "redaction_types": ["financial_card_numbers", "crime_victim_info"]
        }),
    )
    log.info(f"  ✓ Created investigation {INVESTIGATION_CASE_NUMBER}: {new_id}")
    return new_id


# ──────────────────────────────────────────────────────────────────────────────
# Database operations
# ──────────────────────────────────────────────────────────────────────────────

async def get_existing_hashes(pool: asyncpg.Pool) -> set:
    rows = await pool.fetch(
        "SELECT file_hash_sha256 FROM documents WHERE file_hash_sha256 IS NOT NULL"
    )
    return {r["file_hash_sha256"] for r in rows}


async def check_duplicate(pool: asyncpg.Pool, file_hash: str) -> Optional[str]:
    row = await pool.fetchrow(
        "SELECT id::text FROM documents WHERE file_hash_sha256 = $1 LIMIT 1",
        file_hash
    )
    return row["id"] if row else None


async def insert_document(pool: asyncpg.Pool, doc: SystDoc, investigation_id: str) -> str:
    doc_uuid = uuid.UUID(doc.doc_id)
    inv_uuid = uuid.UUID(investigation_id)
    await pool.execute("""
        INSERT INTO documents (
            id, investigation_id, title, document_type, file_path, file_url,
            file_hash_sha256, file_size_bytes, mime_type, description,
            content_text, ocr_text, tags, uploaded_by
        ) VALUES (
            $1, $2, $3, $4, $5, NULL,
            $6, $7, $8, $9,
            $10, $11, $12::text[], $13
        )
        ON CONFLICT (id) DO NOTHING
    """,
        doc_uuid, inv_uuid,
        doc.title[:500], doc.document_type, doc.file_path,
        doc.file_hash_sha256, doc.file_size_bytes, doc.mime_type, doc.description,
        doc.content_text[:2_000_000] if doc.content_text else None,
        doc.ocr_text[:500_000] if doc.ocr_text else None,
        doc.tags,
        "prr_ingestion",
    )
    return doc.doc_id


async def insert_chunks(
    pool: asyncpg.Pool,
    doc_id: str,
    investigation_id: str,
    doc_type: str,
    chunk_tuples: List[Tuple[str, int, int]],
    embeddings: List[List[float]],
):
    """Insert all chunks with embeddings. Skips any chunk that fails to embed."""
    doc_uuid = uuid.UUID(doc_id)
    async with pool.acquire() as conn:
        async with conn.transaction():
            for idx, ((chunk_text_val, start_char, end_char), emb) in enumerate(zip(chunk_tuples, embeddings)):
                token_count = count_tokens_approx(chunk_text_val)
                meta = json.dumps({
                    "document_type": doc_type,
                    "investigation_id": investigation_id,
                    "prr_source": "biia-systemic",
                })
                await conn.execute("""
                    INSERT INTO document_chunks
                        (document_id, chunk_index, chunk_text, embedding,
                         token_count, start_char, end_char, metadata)
                    VALUES (
                        $1, $2, $3, $4::vector,
                        $5, $6, $7, $8::jsonb
                    )
                    ON CONFLICT (document_id, chunk_index) DO NOTHING
                """,
                    doc_uuid, idx, chunk_text_val, vec_literal(emb),
                    token_count, start_char, end_char, meta,
                )


async def create_timeline_event(
    pool: asyncpg.Pool,
    investigation_id: str,
    doc_id: str,
    event_date: str,
    title: str,
    description: str,
    tags: List[str],
    dry_run: bool,
):
    if dry_run:
        return
    inv_uuid = uuid.UUID(investigation_id)
    doc_uuid = uuid.UUID(doc_id)
    # Determine event_date timestamp
    try:
        ts = datetime.fromisoformat(event_date).replace(tzinfo=timezone.utc)
    except Exception:
        ts = datetime.now(timezone.utc)

    # significance: denied=8 (high significance for analysis); granted=5; partial=6
    sig = 5
    if "outcome-denied" in tags:
        sig = 8
    elif "outcome-partial" in tags:
        sig = 6
    elif "outcome-granted" in tags:
        sig = 5

    await pool.execute("""
        INSERT INTO timeline_events (
            investigation_id, event_type, event_date, title,
            description, evidence_ids, significance
        ) VALUES (
            $1, $2, $3, $4,
            $5, $6::uuid[], $7
        )
    """,
        inv_uuid,
        "accommodation_decision",
        ts,
        title[:255],
        description,
        [doc_uuid],
        sig,
    )


async def create_summary_note(
    pool: asyncpg.Pool,
    investigation_id: str,
    content: str,
    dry_run: bool,
):
    if dry_run:
        return
    inv_uuid = uuid.UUID(investigation_id)
    await pool.execute("""
        INSERT INTO notes (investigation_id, title, content, author, note_type)
        VALUES ($1, $2, $3, $4, $5)
    """,
        inv_uuid,
        "PRR Ingestion Complete — Initial Statistics",
        content,
        "ingestion_agent",
        "general",
    )


# ──────────────────────────────────────────────────────────────────────────────
# File hashing + MIME
# ──────────────────────────────────────────────────────────────────────────────

def sha256_file(path: Path) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def get_mime_type(path: Path) -> str:
    ext = path.suffix.lower()
    mime_map = {
        ".pdf":  "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc":  "application/msword",
        ".docm": "application/vnd.ms-word.document.macroEnabled.12",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls":  "application/vnd.ms-excel",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt":  "application/vnd.ms-powerpoint",
        ".msg":  "application/vnd.ms-outlook",
        ".eml":  "message/rfc822",
        ".emlx": "message/rfc822",
        ".txt":  "text/plain",
        ".csv":  "text/csv",
        ".html": "text/html",
        ".htm":  "text/html",
        ".rtf":  "application/rtf",
        ".jpg":  "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png":  "image/png",
        ".tiff": "image/tiff",
        ".tif":  "image/tiff",
        ".bmp":  "image/bmp",
        ".mp3":  "audio/mpeg",
        ".mp4":  "video/mp4",
        ".db":   "application/x-sqlite3",
    }
    return mime_map.get(ext) or mimetypes.guess_type(str(path))[0] or "application/octet-stream"


# ──────────────────────────────────────────────────────────────────────────────
# Validation queries
# ──────────────────────────────────────────────────────────────────────────────

async def run_validation(pool: asyncpg.Pool, investigation_id: str) -> Dict:
    """Run all validation queries and return results as dict."""
    results = {}

    # 1. Total docs
    row = await pool.fetchrow(
        "SELECT COUNT(*) as n FROM documents WHERE investigation_id = $1::uuid",
        investigation_id
    )
    results["total_docs"] = row["n"]

    # 2. By type
    rows = await pool.fetch(
        "SELECT document_type, COUNT(*) as cnt FROM documents WHERE investigation_id = $1::uuid GROUP BY 1 ORDER BY 2 DESC",
        investigation_id
    )
    results["by_type"] = {r["document_type"]: r["cnt"] for r in rows}

    # 3. Missing text
    rows = await pool.fetch(
        """SELECT title, file_path FROM documents
           WHERE investigation_id = $1::uuid
           AND (content_text IS NULL OR LENGTH(content_text) < 10)
           AND (ocr_text IS NULL OR LENGTH(ocr_text) < 10)""",
        investigation_id
    )
    results["missing_text"] = [(r["title"], r["file_path"]) for r in rows]

    # 4. Outcome distribution
    rows = await pool.fetch(
        """SELECT
               CASE
                   WHEN 'outcome-granted' = ANY(tags) THEN 'GRANTED'
                   WHEN 'outcome-denied' = ANY(tags) THEN 'DENIED'
                   WHEN 'outcome-partial' = ANY(tags) THEN 'PARTIAL'
                   ELSE 'NO_OUTCOME'
               END as outcome,
               COUNT(*) as cnt
           FROM documents
           WHERE investigation_id = $1::uuid
           GROUP BY 1 ORDER BY 2 DESC""",
        investigation_id
    )
    results["outcome_dist"] = {r["outcome"]: r["cnt"] for r in rows}

    # 5. Disability type distribution
    rows = await pool.fetch(
        """SELECT tag, COUNT(*) as cnt
           FROM documents, unnest(tags) as tag
           WHERE investigation_id = $1::uuid AND tag LIKE 'disability-%'
           GROUP BY tag ORDER BY cnt DESC""",
        investigation_id
    )
    results["disability_dist"] = {r["tag"]: r["cnt"] for r in rows}

    # 6. Judge distribution
    rows = await pool.fetch(
        """SELECT tag, COUNT(*) as cnt
           FROM documents, unnest(tags) as tag
           WHERE investigation_id = $1::uuid AND tag LIKE 'judge-%'
           GROUP BY tag ORDER BY cnt DESC""",
        investigation_id
    )
    results["judge_dist"] = {r["tag"]: r["cnt"] for r in rows}

    # 7. Accommodation type distribution
    rows = await pool.fetch(
        """SELECT tag, COUNT(*) as cnt
           FROM documents, unnest(tags) as tag
           WHERE investigation_id = $1::uuid AND tag LIKE 'accom-%'
           GROUP BY tag ORDER BY cnt DESC""",
        investigation_id
    )
    results["accom_dist"] = {r["tag"]: r["cnt"] for r in rows}

    # 8. Grant/deny by disability type
    rows = await pool.fetch(
        """SELECT
               disability_tag,
               SUM(CASE WHEN 'outcome-granted' = ANY(d.tags) THEN 1 ELSE 0 END) as granted,
               SUM(CASE WHEN 'outcome-denied' = ANY(d.tags) THEN 1 ELSE 0 END) as denied,
               SUM(CASE WHEN 'outcome-partial' = ANY(d.tags) THEN 1 ELSE 0 END) as partial,
               COUNT(*) as total_decisions
           FROM documents d, unnest(d.tags) as disability_tag
           WHERE d.investigation_id = $1::uuid
           AND disability_tag LIKE 'disability-%'
           AND ('outcome-granted' = ANY(d.tags) OR 'outcome-denied' = ANY(d.tags) OR 'outcome-partial' = ANY(d.tags))
           GROUP BY disability_tag ORDER BY total_decisions DESC""",
        investigation_id
    )
    results["grant_deny_by_disability"] = [
        {
            "disability": r["disability_tag"],
            "granted": r["granted"],
            "denied": r["denied"],
            "partial": r["partial"],
            "total": r["total_decisions"],
        }
        for r in rows
    ]

    # 9. Grant/deny by judge
    rows = await pool.fetch(
        """SELECT
               judge_tag,
               SUM(CASE WHEN 'outcome-granted' = ANY(d.tags) THEN 1 ELSE 0 END) as granted,
               SUM(CASE WHEN 'outcome-denied' = ANY(d.tags) THEN 1 ELSE 0 END) as denied,
               SUM(CASE WHEN 'outcome-partial' = ANY(d.tags) THEN 1 ELSE 0 END) as partial,
               COUNT(*) as total_decisions
           FROM documents d, unnest(d.tags) as judge_tag
           WHERE d.investigation_id = $1::uuid
           AND judge_tag LIKE 'judge-%'
           AND ('outcome-granted' = ANY(d.tags) OR 'outcome-denied' = ANY(d.tags) OR 'outcome-partial' = ANY(d.tags))
           GROUP BY judge_tag ORDER BY total_decisions DESC""",
        investigation_id
    )
    results["grant_deny_by_judge"] = [
        {
            "judge": r["judge_tag"],
            "granted": r["granted"],
            "denied": r["denied"],
            "partial": r["partial"],
            "total": r["total_decisions"],
        }
        for r in rows
    ]

    # 10. Timeline events
    row = await pool.fetchrow(
        "SELECT COUNT(*) as n FROM timeline_events WHERE investigation_id = $1::uuid",
        investigation_id
    )
    results["timeline_events"] = row["n"]

    # 11. Chunks and embeddings
    row = await pool.fetchrow(
        """SELECT
               COUNT(DISTINCT dc.document_id) as docs_with_chunks,
               COUNT(dc.id) as total_chunks,
               SUM(CASE WHEN dc.embedding IS NOT NULL THEN 1 ELSE 0 END) as chunks_with_embeddings,
               SUM(CASE WHEN dc.embedding IS NULL THEN 1 ELSE 0 END) as chunks_missing_embeddings
           FROM document_chunks dc
           JOIN documents d ON dc.document_id = d.id
           WHERE d.investigation_id = $1::uuid""",
        investigation_id
    )
    results["chunks_stats"] = dict(row) if row else {}

    # 12. Tag frequency (top 30)
    rows = await pool.fetch(
        """SELECT tag, COUNT(*) as cnt
           FROM documents, unnest(tags) as tag
           WHERE investigation_id = $1::uuid
           GROUP BY tag ORDER BY cnt DESC LIMIT 30""",
        investigation_id
    )
    results["tag_freq"] = {r["tag"]: r["cnt"] for r in rows}

    return results


# ──────────────────────────────────────────────────────────────────────────────
# Report generation
# ──────────────────────────────────────────────────────────────────────────────

def generate_report(
    stats: Dict,
    validation: Dict,
    investigation_id: str,
    run_started: datetime,
) -> str:
    run_ended = datetime.now(timezone.utc)
    duration = (run_ended - run_started).total_seconds()

    lines = [
        "# BIIA Systemic PRR Ingestion Report",
        "",
        f"**Investigation**: {INVESTIGATION_CASE_NUMBER}  ",
        f"**Investigation UUID**: `{investigation_id}`  ",
        f"**Run started**: {run_started.strftime('%Y-%m-%d %H:%M:%S UTC')}  ",
        f"**Run completed**: {run_ended.strftime('%Y-%m-%d %H:%M:%S UTC')}  ",
        f"**Duration**: {int(duration // 60)}m {int(duration % 60)}s  ",
        "",
        "---",
        "",
        "## 1. Ingestion Summary",
        "",
        f"| Metric | Count |",
        f"|--------|-------|",
        f"| Files found | {stats.get('files_found', 0):,} |",
        f"| Files skipped (system/temp) | {stats.get('files_skipped', 0):,} |",
        f"| Files attempted | {stats.get('files_attempted', 0):,} |",
        f"| **Documents inserted** | **{stats.get('docs_inserted', 0):,}** |",
        f"| Duplicates skipped | {stats.get('duplicates', 0):,} |",
        f"| Extraction failures | {stats.get('extraction_failures', 0):,} |",
        f"| Total chunks created | {stats.get('total_chunks', 0):,} |",
        f"| Total embeddings generated | {stats.get('total_embeddings', 0):,} |",
        f"| Timeline events created | {validation.get('timeline_events', 0):,} |",
        f"| OpenAI tokens used (ingestion) | {stats.get('tokens_used', 0):,} |",
        "",
        "---",
        "",
        "## 2. Document Type Distribution",
        "",
        "| Document Type | Count |",
        "|---------------|-------|",
    ]
    for dtype, cnt in sorted(validation.get("by_type", {}).items(), key=lambda x: -x[1]):
        lines.append(f"| `{dtype}` | {cnt:,} |")

    lines += [
        "",
        "---",
        "",
        "## 3. Accommodation Outcome Distribution",
        "",
        "| Outcome | Count |",
        "|---------|-------|",
    ]
    for outcome, cnt in sorted(validation.get("outcome_dist", {}).items(), key=lambda x: -x[1]):
        lines.append(f"| {outcome} | {cnt:,} |")

    lines += [
        "",
        "---",
        "",
        "## 4. ★ GRANT/DENY RATE BY DISABILITY TYPE ★",
        "",
        "_This is the primary systemic analysis metric._",
        "",
        "| Disability Type | Granted | Denied | Partial | Total | Grant % | Deny % |",
        "|----------------|---------|--------|---------|-------|---------|--------|",
    ]
    for row in validation.get("grant_deny_by_disability", []):
        total = row["total"] or 1
        grant_pct = f"{100 * row['granted'] / total:.0f}%"
        deny_pct  = f"{100 * row['denied'] / total:.0f}%"
        lines.append(
            f"| `{row['disability']}` | {row['granted']} | {row['denied']} | {row['partial']} "
            f"| {row['total']} | {grant_pct} | {deny_pct} |"
        )

    lines += [
        "",
        "---",
        "",
        "## 5. ★ GRANT/DENY RATE BY JUDGE ★",
        "",
        "| Judge | Granted | Denied | Partial | Total | Grant % | Deny % |",
        "|-------|---------|--------|---------|-------|---------|--------|",
    ]
    for row in validation.get("grant_deny_by_judge", []):
        total = row["total"] or 1
        grant_pct = f"{100 * row['granted'] / total:.0f}%"
        deny_pct  = f"{100 * row['denied'] / total:.0f}%"
        lines.append(
            f"| `{row['judge']}` | {row['granted']} | {row['denied']} | {row['partial']} "
            f"| {row['total']} | {grant_pct} | {deny_pct} |"
        )

    lines += [
        "",
        "---",
        "",
        "## 6. Disability Type Distribution",
        "",
        "| Disability Tag | Documents |",
        "|----------------|-----------|",
    ]
    for tag, cnt in sorted(validation.get("disability_dist", {}).items(), key=lambda x: -x[1]):
        lines.append(f"| `{tag}` | {cnt:,} |")

    lines += [
        "",
        "---",
        "",
        "## 7. Judge Distribution",
        "",
        "| Judge Tag | Documents |",
        "|-----------|-----------|",
    ]
    for tag, cnt in sorted(validation.get("judge_dist", {}).items(), key=lambda x: -x[1]):
        lines.append(f"| `{tag}` | {cnt:,} |")

    lines += [
        "",
        "---",
        "",
        "## 8. Accommodation Type Distribution",
        "",
        "| Accommodation Type | Documents |",
        "|--------------------|-----------|",
    ]
    for tag, cnt in sorted(validation.get("accom_dist", {}).items(), key=lambda x: -x[1]):
        lines.append(f"| `{tag}` | {cnt:,} |")

    lines += [
        "",
        "---",
        "",
        "## 9. Tag Frequency (Top 30)",
        "",
        "| Tag | Count |",
        "|-----|-------|",
    ]
    for tag, cnt in list(validation.get("tag_freq", {}).items())[:30]:
        lines.append(f"| `{tag}` | {cnt:,} |")

    lines += [
        "",
        "---",
        "",
        "## 10. Chunks & Embeddings",
        "",
    ]
    cs = validation.get("chunks_stats", {})
    lines += [
        f"| Metric | Value |",
        f"|--------|-------|",
        f"| Documents with chunks | {cs.get('docs_with_chunks', 0):,} |",
        f"| Total chunks | {cs.get('total_chunks', 0):,} |",
        f"| Chunks with embeddings | {cs.get('chunks_with_embeddings', 0):,} |",
        f"| Chunks missing embeddings | {cs.get('chunks_missing_embeddings', 0):,} |",
    ]

    lines += [
        "",
        "---",
        "",
        "## 11. Nelson-Related Documents Found",
        "",
        f"Documents tagged `nelson-related`: **{stats.get('nelson_related_count', 0):,}**",
        "",
        "---",
        "",
        "## 12. Duplicates Found",
        "",
    ]
    dupes = stats.get("duplicate_list", [])
    if dupes:
        lines += [
            "| File | Existing Hash |",
            "|------|---------------|",
        ]
        for fname, fhash in dupes[:50]:
            lines.append(f"| `{fname}` | `{fhash[:12]}…` |")
        if len(dupes) > 50:
            lines.append(f"_…and {len(dupes)-50} more_")
    else:
        lines.append("_No duplicates found._")

    lines += [
        "",
        "---",
        "",
        "## 13. Text Extraction Failures",
        "",
    ]
    failures = validation.get("missing_text", [])
    if failures:
        lines += [
            "| Title | File Path |",
            "|-------|-----------|",
        ]
        for title, fpath in failures[:50]:
            lines.append(f"| {title} | `{fpath}` |")
        if len(failures) > 50:
            lines.append(f"_…and {len(failures)-50} more_")
    else:
        lines.append("_No extraction failures (all documents have content text)._")

    lines += [
        "",
        "---",
        "",
        "## 14. Embedding Upgrade",
        "",
    ]
    eu = stats.get("embedding_upgrade", {})
    if eu.get("skipped"):
        lines.append("_Embeddings were already at 3072 dimensions — upgrade skipped._")
    else:
        lines += [
            f"- Document chunks re-embedded: {eu.get('chunks_processed', 0):,}",
            f"- Entity embeddings re-embedded: {eu.get('entity_embs_processed', 0):,}",
            f"- Upgrade complete: both tables now use `{EMBEDDING_MODEL}` (3072 dims)",
        ]

    lines += [
        "",
        "---",
        "",
        "_Report generated by ingest_biia_systemic_prr.py_",
    ]

    return "\n".join(lines)


# ──────────────────────────────────────────────────────────────────────────────
# Main ingestion loop
# ──────────────────────────────────────────────────────────────────────────────

async def main():
    parser = argparse.ArgumentParser(description="BIIA Systemic PRR Ingestion")
    parser.add_argument("--dry-run", action="store_true",
                        help="Extract and classify without writing to DB")
    parser.add_argument("--resume", action="store_true",
                        help="Skip files whose hash already exists in DB")
    parser.add_argument("--skip-embed-upgrade", action="store_true",
                        help="Skip the embedding dimension upgrade step")
    args = parser.parse_args()

    dry_run = args.dry_run
    resume  = args.resume

    if dry_run:
        log.info("=== DRY RUN MODE — no DB writes ===")

    run_started = datetime.now(timezone.utc)

    # ── Connect to osint-db ─────────────────────────────────────────────────
    log.info("Connecting to osint-db…")
    pool = await asyncpg.create_pool(DB_URL, min_size=2, max_size=10)
    embedder = EmbeddingClient()

    stats = {
        "files_found": 0,
        "files_skipped": 0,
        "files_attempted": 0,
        "docs_inserted": 0,
        "duplicates": 0,
        "extraction_failures": 0,
        "total_chunks": 0,
        "total_embeddings": 0,
        "tokens_used": 0,
        "nelson_related_count": 0,
        "duplicate_list": [],
        "embedding_upgrade": {},
    }

    try:
        # ── Step 0: Upgrade embeddings ──────────────────────────────────────
        if not args.skip_embed_upgrade:
            eu_stats = await upgrade_embeddings(pool, embedder, dry_run)
            stats["embedding_upgrade"] = eu_stats
        else:
            log.info("Skipping embedding upgrade (--skip-embed-upgrade)")
            stats["embedding_upgrade"] = {"skipped": True}

        # ── Step 1: Create investigation ────────────────────────────────────
        log.info("=== STEP 1: Get/create investigation ===")
        investigation_id = await get_or_create_investigation(pool, dry_run)

        # Load existing hashes for dedup
        existing_hashes: set = set()
        if resume and not dry_run:
            existing_hashes = await get_existing_hashes(pool)
            log.info(f"  Loaded {len(existing_hashes):,} existing hashes for dedup")

        # ── Step 2: Walk directory ───────────────────────────────────────────
        log.info(f"=== STEP 2: Walking {SOURCE_ROOT} ===")
        all_files: List[Path] = []
        for root, dirs, files in os.walk(SOURCE_ROOT):
            # Skip __MACOSX dirs
            dirs[:] = [d for d in dirs if d != "__MACOSX" and not d.startswith(".")]
            for fname in files:
                all_files.append(Path(root) / fname)

        stats["files_found"] = len(all_files)
        log.info(f"  Found {len(all_files):,} total files")

        # Filter skippable files
        process_files = []
        for fp in all_files:
            if should_skip(fp):
                stats["files_skipped"] += 1
            else:
                process_files.append(fp)

        log.info(f"  Skipping {stats['files_skipped']:,} system/temp files")
        log.info(f"  Processing {len(process_files):,} files")

        # ── Steps 3–8: Process each file ────────────────────────────────────
        for i, file_path in enumerate(process_files, 1):
            stats["files_attempted"] += 1
            rel_name = str(file_path.relative_to(SOURCE_ROOT.parent.parent))

            if i % 50 == 0 or i == 1:
                log.info(f"  [{i}/{len(process_files)}] {rel_name[:80]}")

            try:
                # Compute hash
                file_hash = sha256_file(file_path)
                file_size = file_path.stat().st_size

                # Dedup check
                if resume and file_hash in existing_hashes:
                    stats["duplicates"] += 1
                    stats["duplicate_list"].append((file_path.name, file_hash))
                    log.debug(f"    DUPLICATE (hash match): {file_path.name}")
                    continue
                if not resume and not dry_run:
                    existing_id = await check_duplicate(pool, file_hash)
                    if existing_id:
                        stats["duplicates"] += 1
                        stats["duplicate_list"].append((file_path.name, file_hash))
                        log.debug(f"    DUPLICATE (hash match): {file_path.name}")
                        continue

                # Extract text
                content_text, method, ocr_text = extract_content(file_path)

                if not content_text and file_path.suffix.lower() in SKIP_CONTENT_EXTS:
                    stats["files_skipped"] += 1
                    continue

                if not content_text or len(content_text) < 5:
                    stats["extraction_failures"] += 1
                    log.debug(f"    EXTRACTION FAILURE: {file_path.name} (method={method})")

                # Date extraction
                ext_lower = file_path.suffix.lower()
                doc_date = None
                if ext_lower == ".msg":
                    doc_date = extract_msg_date(file_path)
                if not doc_date:
                    doc_date = extract_doc_date(content_text or "", file_path.name)
                if not doc_date:
                    _, installment_date = detect_installment(file_path)
                    doc_date = installment_date

                # Classify
                doc_type = classify_document(file_path, content_text or "")

                # Build tags + extract metadata
                tags, judge_name, outcome = build_tags(file_path, content_text or "", doc_type)

                # Build title + description
                title = build_title(file_path, content_text or "")
                description = build_description(file_path, content_text or "",
                                                doc_type, tags, judge_name, outcome)

                # Create document
                doc_id  = str(uuid.uuid4())
                mime    = get_mime_type(file_path)

                doc = SystDoc(
                    doc_id=doc_id,
                    file_path=str(file_path),
                    filename=file_path.name,
                    title=title,
                    document_type=doc_type,
                    file_hash_sha256=file_hash,
                    file_size_bytes=file_size,
                    mime_type=mime,
                    description=description,
                    content_text=content_text or "",
                    ocr_text=ocr_text,
                    tags=tags,
                    extraction_method=method,
                    doc_date=doc_date,
                )

                if "nelson-related" in tags:
                    stats["nelson_related_count"] += 1

                if not dry_run:
                    await insert_document(pool, doc, investigation_id)
                    existing_hashes.add(file_hash)  # prevent intra-run dupes

                    # Chunk + embed (skip if no content)
                    if content_text and len(content_text) >= 20:
                        chunk_tuples = chunk_text(content_text)
                        if chunk_tuples:
                            chunk_texts_only = [ct for ct, _, _ in chunk_tuples]
                            embeddings = await embedder.embed_batched(chunk_texts_only)
                            await insert_chunks(
                                pool, doc_id, investigation_id, doc_type,
                                chunk_tuples, embeddings
                            )
                            stats["total_chunks"] += len(chunk_tuples)
                            stats["total_embeddings"] += len(embeddings)

                    # Timeline event for accommodation decisions
                    if doc_type == "accommodation_order" and outcome in ("granted", "denied", "partial"):
                        outcome_label = {
                            "granted": "GRANTED",
                            "denied":  "DENIED",
                            "partial": "PARTIALLY GRANTED",
                        }.get(outcome, outcome.upper())
                        judge_part = f" — {judge_name}" if judge_name else ""
                        event_title = f"BIIA {outcome_label} accommodation{judge_part}"
                        if len(tags) > 0:
                            accomm_tags = [t for t in tags if t.startswith("accom-")]
                            if accomm_tags:
                                accomm_str = accomm_tags[0].replace("accom-", "")
                                event_title = f"BIIA {outcome_label} accommodation ({accomm_str}){judge_part}"
                        event_date = doc_date or datetime.now(timezone.utc).strftime("%Y-%m-%d")
                        await create_timeline_event(
                            pool, investigation_id, doc_id,
                            event_date, event_title, description, tags, dry_run
                        )

                stats["docs_inserted"] += 1

            except Exception as e:
                log.error(f"    ERROR processing {file_path.name}: {e}", exc_info=True)
                stats["extraction_failures"] += 1

        stats["tokens_used"] = embedder.total_tokens

        # ── Step 9: Run validation ───────────────────────────────────────────
        log.info("=== STEP 9: Running validation queries ===")
        if not dry_run:
            validation = await run_validation(pool, investigation_id)
            log.info(f"  Total docs in DB: {validation.get('total_docs', 0):,}")
            log.info(f"  Outcome dist: {validation.get('outcome_dist', {})}")
        else:
            validation = {
                "total_docs": stats["docs_inserted"],
                "by_type": {}, "missing_text": [], "outcome_dist": {},
                "disability_dist": {}, "judge_dist": {}, "accom_dist": {},
                "grant_deny_by_disability": [], "grant_deny_by_judge": [],
                "timeline_events": 0, "chunks_stats": {}, "tag_freq": {},
            }

        # ── Step 10: Summary note ────────────────────────────────────────────
        report_md = generate_report(stats, validation, investigation_id, run_started)

        if not dry_run:
            note_content = (
                f"## Ingestion Summary\n\n"
                f"- Total documents: {validation.get('total_docs', 0):,}\n"
                f"- Total chunks: {stats['total_chunks']:,}\n"
                f"- Total embeddings: {stats['total_embeddings']:,}\n"
                f"- Timeline events: {validation.get('timeline_events', 0):,}\n"
                f"- Duplicates skipped: {stats['duplicates']:,}\n"
                f"- Extraction failures: {stats['extraction_failures']:,}\n"
                f"- Nelson-related docs: {stats['nelson_related_count']:,}\n"
                f"- OpenAI tokens used: {stats['tokens_used']:,}\n\n"
                f"### Outcome Distribution\n"
            )
            for outcome, cnt in validation.get("outcome_dist", {}).items():
                note_content += f"- {outcome}: {cnt:,}\n"
            note_content += "\n### Top Tags\n"
            for tag, cnt in list(validation.get("tag_freq", {}).items())[:15]:
                note_content += f"- `{tag}`: {cnt:,}\n"

            await create_summary_note(pool, investigation_id, note_content, dry_run)

        # ── Step 11: Write report ────────────────────────────────────────────
        log.info(f"=== STEP 11: Writing report to {REPORT_PATH} ===")
        REPORT_PATH.write_text(report_md, encoding="utf-8")
        log.info(f"  ✓ Report written")

        # Final summary
        log.info("=== INGESTION COMPLETE ===")
        log.info(f"  Files found:          {stats['files_found']:,}")
        log.info(f"  Files attempted:      {stats['files_attempted']:,}")
        log.info(f"  Documents inserted:   {stats['docs_inserted']:,}")
        log.info(f"  Duplicates skipped:   {stats['duplicates']:,}")
        log.info(f"  Extraction failures:  {stats['extraction_failures']:,}")
        log.info(f"  Total chunks:         {stats['total_chunks']:,}")
        log.info(f"  Total embeddings:     {stats['total_embeddings']:,}")
        log.info(f"  Timeline events:      {validation.get('timeline_events', 0):,}")
        log.info(f"  OpenAI tokens used:   {stats['tokens_used']:,}")

    finally:
        await embedder.close()
        await pool.close()


if __name__ == "__main__":
    asyncio.run(main())
