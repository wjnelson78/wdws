#!/usr/bin/env python3
"""
BIIA Nelson PRR Ingestion Script
Ingests public records from BIIA docket 25-18153 / claim SX-21824 into wdws DB.

Usage:
    python3 ingest_biia_nelson.py [--dry-run] [--resume]

    --dry-run   Extract and classify without writing to DB
    --resume    Skip files whose content_hash already exists in core.documents
                as a biia-prr-nelson source (default: always dedup by hash)
"""

import argparse
import asyncio
import email as email_lib
from email import policy as email_policy
import hashlib
import json
import logging
import os
import re
import subprocess
import sys
sys.path.insert(0, "/opt/wdws")
import time
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

import asyncpg
import httpx
import pytesseract
from PIL import Image

from embedding_service import (
    embed_texts_sync, embed_query_sync,
    _vec_literal as _embedding_vec_literal,
    EMBEDDING_DIMENSIONS, EMBEDDING_MODEL,
)
from contextual_retrieval import generate_context_sync, enrich_chunks

# ──────────────────────────────────────────────────────────────────────────────
# Configuration
# ──────────────────────────────────────────────────────────────────────────────

DATABASE_URL = os.environ["DATABASE_URL"]
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# EMBEDDING_MODEL and EMBEDDING_DIMENSIONS imported from embedding_service
# (BGE-M3 local model, 1024 dimensions)
EMBEDDING_DIMS   = EMBEDDING_DIMENSIONS
EMBEDDING_MODEL_ID = 1
EMBEDDING_BATCH  = 50   # chunks per embed request

SOURCE_ROOT = Path(
    "/opt/wdws/data/dropbox/biia/"
    "09.18.2025 Nelson Request 2518153 Accommodation Records 6.1.25-9.17.25/"
    "09.18.2025 Nelson Request 2518153 Accommodation Records 6.1.25-9.17.25"
)
REPORT_PATH = Path(
    "/opt/wdws/data/dropbox/biia/INGESTION_REPORT_NELSON_PRR.md"
)

# Token-approximate chunking
CHUNK_TARGET_CHARS = 2500   # ~600 tokens at 4 chars/token
CHUNK_OVERLAP_CHARS = 200

OCR_DPI = 200

# SKIP patterns
SKIP_NAMES = {".ds_store", "thumbs.db", "desktop.ini"}
SKIP_PREFIXES = ("._", "__macosx")

# Files that are actually installment subfolders' macOS AppleDouble artifacts
SKIP_EXTENSIONS = {".ds_store"}

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    datefmt="%H:%M:%S",
)
log = logging.getLogger("biia_ingest")

if not OPENAI_API_KEY:
    raise RuntimeError("OPENAI_API_KEY environment variable is required")

# ──────────────────────────────────────────────────────────────────────────────
# Data model
# ──────────────────────────────────────────────────────────────────────────────

@dataclass
class BiiaDocument:
    doc_id: str
    source_path: str
    filename: str
    document_type: str
    title: str
    content_hash: str
    full_content: str
    metadata: Dict
    chunks: List[str] = field(default_factory=list)
    extraction_method: str = "unknown"
    is_duplicate: bool = False           # hash already existed in DB
    existing_doc_id: Optional[str] = None   # UUID of existing doc if duplicate


# ──────────────────────────────────────────────────────────────────────────────
# Installment detection
# ──────────────────────────────────────────────────────────────────────────────

INSTALLMENT_MAP = {
    "1st": 1, "2nd": 2, "3rd": 3, "4th": 4,
    "5th": 5, "6th": 6, "7th": 7,
}

def detect_installment(file_path: Path) -> Optional[int]:
    """Determine installment number from parent directory name."""
    for parent in file_path.parents:
        name_lower = parent.name.lower()
        for prefix, num in INSTALLMENT_MAP.items():
            if name_lower.startswith(prefix):
                return num
    return None


# ──────────────────────────────────────────────────────────────────────────────
# Text extraction
# ──────────────────────────────────────────────────────────────────────────────

def extract_pdf(pdf_path: Path) -> Tuple[str, str]:
    """
    Returns (text, method).
    Tries pdftotext first; falls back to pymupdf text layer; then OCR.
    """
    # 1. pdftotext
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", "-enc", "UTF-8", str(pdf_path), "-"],
            capture_output=True, text=True, timeout=60
        )
        text = result.stdout.strip()
        if text and len(text) > 50 and not _is_garbage(text):
            return text, "pdftotext"
    except Exception as e:
        log.debug(f"pdftotext failed for {pdf_path.name}: {e}")

    # 2. PyMuPDF text layer
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        parts = []
        for page_num, page in enumerate(doc, 1):
            t = page.get_text("text")
            if t and t.strip():
                parts.append(f"[Page {page_num}]\n{t}")
        doc.close()
        text = "\n\n".join(parts).strip()
        if text and len(text) > 50 and not _is_garbage(text):
            return text, "pymupdf"
    except Exception as e:
        log.debug(f"PyMuPDF text layer failed for {pdf_path.name}: {e}")

    # 3. OCR via PyMuPDF render + tesseract
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        parts = []
        for page_num, page in enumerate(doc, 1):
            mat = fitz.Matrix(OCR_DPI / 72, OCR_DPI / 72)
            pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            t = pytesseract.image_to_string(img)
            if t and t.strip():
                parts.append(f"[Page {page_num}]\n{t}")
        doc.close()
        text = "\n\n".join(parts).strip()
        if text:
            return text, "tesseract_ocr"
    except Exception as e:
        log.warning(f"OCR failed for {pdf_path.name}: {e}")

    return "", "failed"


def _is_garbage(text: str) -> bool:
    """Heuristic: if >40% non-ASCII, treat as garbled."""
    if not text:
        return True
    sample = text[:2000]
    non_ascii = sum(1 for c in sample if ord(c) > 127)
    return (non_ascii / max(len(sample), 1)) > 0.40


def extract_docx(path: Path) -> Tuple[str, str]:
    try:
        import docx as docx_lib
        doc = docx_lib.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        paragraphs.append(cell.text.strip())
        return "\n".join(paragraphs), "docx_parser"
    except Exception as e:
        log.warning(f"DOCX extraction failed for {path.name}: {e}")
        return "", "failed"


def extract_msg(path: Path) -> Tuple[str, str]:
    """Extract text from Outlook .msg file."""
    try:
        import extract_msg as em_lib
        msg = em_lib.Message(str(path))
        parts = []
        if msg.subject:
            parts.append(f"Subject: {msg.subject}")
        if msg.sender:
            parts.append(f"From: {msg.sender}")
        if msg.to:
            parts.append(f"To: {msg.to}")
        if msg.cc:
            parts.append(f"CC: {msg.cc}")
        if msg.date:
            parts.append(f"Date: {msg.date}")
        parts.append("")
        body = msg.body or ""
        if body:
            parts.append(body)
        return "\n".join(parts), "extract_msg"
    except Exception as e:
        log.warning(f"MSG extraction failed for {path.name}: {e}")
        return "", "failed"


def extract_eml(path: Path) -> Tuple[str, str]:
    try:
        with open(path, "rb") as f:
            msg = email_lib.message_from_binary_file(f, policy=email_policy.default)
        parts = []
        for h in ("Subject", "From", "To", "CC", "Date"):
            val = msg.get(h, "")
            if val:
                parts.append(f"{h}: {val}")
        parts.append("")
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                cd = str(part.get("Content-Disposition", ""))
                if "attachment" in cd.lower():
                    continue
                if ct == "text/plain":
                    try:
                        body += part.get_content()
                    except Exception:
                        raw = part.get_payload(decode=True)
                        if raw:
                            body += raw.decode("utf-8", errors="ignore")
        else:
            try:
                body = msg.get_content()
            except Exception:
                raw = msg.get_payload(decode=True)
                if raw:
                    body = raw.decode("utf-8", errors="ignore")
        parts.append(body)
        return "\n".join(parts), "python_email"
    except Exception as e:
        log.warning(f"EML extraction failed for {path.name}: {e}")
        return "", "failed"


def extract_image(path: Path) -> Tuple[str, str]:
    try:
        img = Image.open(str(path))
        text = pytesseract.image_to_string(img)
        return text.strip(), "tesseract_ocr"
    except Exception as e:
        log.warning(f"Image OCR failed for {path.name}: {e}")
        return "", "failed"


def extract_text_file(path: Path) -> Tuple[str, str]:
    try:
        return path.read_text(encoding="utf-8", errors="ignore"), "text_file"
    except Exception as e:
        log.warning(f"Text file read failed for {path.name}: {e}")
        return "", "failed"


def extract_content(path: Path) -> Tuple[str, str]:
    """Dispatch to appropriate extractor. Returns (text, method)."""
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path)
    elif ext == ".docx":
        return extract_docx(path)
    elif ext == ".msg":
        return extract_msg(path)
    elif ext in (".eml", ".emlx"):
        return extract_eml(path)
    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"):
        return extract_image(path)
    elif ext in (".txt", ".csv", ".html", ".htm"):
        return extract_text_file(path)
    elif ext == ".xlsx":
        return extract_xlsx(path)
    elif ext in (".pptx", ".ppt"):
        return extract_pptx(path)
    else:
        log.info(f"  Unknown extension {ext} for {path.name}, trying text read")
        return extract_text_file(path)


def extract_pptx(path: Path) -> Tuple[str, str]:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            parts.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n\n".join(parts), "python_pptx"
    except Exception as e:
        log.warning(f"PPTX extraction failed for {path.name}: {e}")
        return "", "failed"


def extract_xlsx(path: Path) -> Tuple[str, str]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        rows = []
        for sheet in wb.worksheets:
            rows.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                line = "\t".join(str(v) for v in row if v is not None)
                if line.strip():
                    rows.append(line)
        return "\n".join(rows), "openpyxl"
    except ImportError:
        # Fall back to csv-like read
        return extract_text_file(path)
    except Exception as e:
        log.warning(f"XLSX extraction failed for {path.name}: {e}")
        return "", "failed"


# ──────────────────────────────────────────────────────────────────────────────
# Classification
# ──────────────────────────────────────────────────────────────────────────────

# Judge name extraction
JUDGE_PATTERNS = [
    (re.compile(r'Timothy\s+M\.?\s+Blood', re.I), 'Timothy M. Blood'),
    (re.compile(r'Tim\s+Blood', re.I), 'Timothy M. Blood'),
    (re.compile(r'Judge\s+Blood', re.I), 'Timothy M. Blood'),
    (re.compile(r'Ann\s+(?:M\.?\s+)?Dodge', re.I), 'Ann Dodge'),
    (re.compile(r'Judge\s+Dodge', re.I), 'Ann Dodge'),
    (re.compile(r'Brian\s+(?:C\.?\s+)?Watkins', re.I), 'Brian Watkins'),
    (re.compile(r'Judge\s+Watkins', re.I), 'Brian Watkins'),
    (re.compile(r'Janice\s+(?:K\.?\s+)?Rosen', re.I), 'Janice Rosen'),
    (re.compile(r'Board\s+Member\s+Rosen', re.I), 'Janice Rosen'),
]

ACCOMM_TYPE_PATTERNS = {
    "email_filing":        re.compile(r'email\s+fil(ing|e)', re.I),
    "email_service":       re.compile(r'email\s+service|service\s+by\s+email', re.I),
    "remote_hearing":      re.compile(r'remote\s+hear|telephon(ic|e)\s+hear|video\s+hear', re.I),
    "extended_time":       re.compile(r'extend(ed)?\s+time|additional\s+time|time\s+extension', re.I),
    "counsel_appointment": re.compile(r'appoint(ment)?\s+of\s+counsel|appointed\s+counsel', re.I),
    "gal_appointment":     re.compile(r'guardian\s+ad\s+litem|GAL\s+appoint', re.I),
    "assistive_technology":re.compile(r'assistive\s+tech|AI\s+assist|Athena', re.I),
    "stay_of_proceedings": re.compile(r'stay\s+of\s+proceed|stay\s+the\s+case', re.I),
    "physical_access":     re.compile(r'physical\s+access|wheelchair|parking', re.I),
    "interpreter":         re.compile(r'interpreter|translation|language\s+access', re.I),
}

ACCOMM_ACTION_PATTERNS = {
    "granted":          re.compile(r'\b(grant(ed|ing|s)?|allow(ed|ing)?)\b.{0,80}(accommodation|request)', re.I),
    "denied":           re.compile(r'\b(den(ied|y|ying)|declin(ed|ing)|refus(ed|ing))\b.{0,80}(accommodation|request)', re.I),
    "partially_granted":re.compile(r'partial(ly)?\s+(grant(ed)?|allow(ed)?)', re.I),
    "deferred":         re.compile(r'defer(red)?|hold.?in\s+abey|tabled', re.I),
}

FILE_BASED_DOCTYPE = {
    # High-specificity mailroom signals checked early
    "mailing details":   "mailroom_record",
    "bams":              "mailroom_record",
    "certified mail":    "mailroom_record",
    "usps":              "mailroom_record",
    "return receipt":    "mailroom_record",
    # Court / legal
    "appeal":            "court_filing",
    "filing":            "court_filing",
    "motion":            "court_filing",
    "petition":          "court_filing",
    # Orders
    "order":             "accommodation_order",
    "ord_":              "accommodation_order",
    # Policy
    "dist list":         "accommodation_policy",
    "distribution list": "accommodation_policy",
    "triage info":       "accommodation_policy",
    "policy":            "accommodation_policy",
    # Communications
    "letter":            "external_communication",
    "cor_":              "external_communication",
    "email":             "internal_communication",
    "log":               "internal_communication",
    "case log":          "internal_communication",
    "active cases":      "internal_communication",
    # Phone
    "voicemail":         "phone_record",
    "phone":             "phone_record",
    "cdr":               "phone_record",
    # Drafts last (low specificity)
    "draft":             "draft_document",
}

FOLDER_DOCTYPE = {
    "mailroom records":   "mailroom_record",
    "mailroom":           "mailroom_record",
    "drafts and orders":  "draft_document",
    "emailed correspondence": "external_communication",
    "dist lists":         "accommodation_policy",
    "policies etc":       "accommodation_policy",
    "outlook":            "internal_communication",
    "voicemail":          "phone_record",
    "communications with ago": "internal_communication",
    "folders-drives":     "internal_communication",
}

CONTENT_DOCTYPE_PATTERNS = [
    (re.compile(r'ORDER\s+(ON|GRANTING|DENYING|RE:|REGARDING)\s+ACCOMMODATION', re.I), "accommodation_order"),
    (re.compile(r'ORDER\s+ON\s+.+ACCOMMODATION', re.I), "accommodation_order"),
    (re.compile(r'ORDER\s+STAYING\s+ACTION', re.I), "accommodation_order"),
    (re.compile(r'General Order.{0,50}Nelson', re.I), "accommodation_order"),
    (re.compile(r'(?:BIIA|Board).{0,30}ACCOMMODATION\s+REQUEST', re.I), "accommodation_request"),
    (re.compile(r'I\s+(?:hereby\s+)?request\s+(an?\s+)?accommodation', re.I), "accommodation_request"),
    (re.compile(r'ADA\s+ACCOMMODATION\s+REQUEST\s+FORM', re.I), "accommodation_request"),
    (re.compile(r'(?:BIIA|Board).{0,30}(?:POLICY|PROCEDURE|GUIDELINES?)\s+(?:FOR|ON|RE)\s+ACCOMMODATION', re.I), "accommodation_policy"),
    (re.compile(r'CERTIFIED\s+MAIL(?:ING)?', re.I), "mailroom_record"),
    (re.compile(r'USPS\s+(?:Return\s+Receipt|Tracking)', re.I), "mailroom_record"),
    (re.compile(r'BAMS\s+Mailing', re.I), "mailroom_record"),
    (re.compile(r'voicemail', re.I), "phone_record"),
]


def classify_document(path: Path, text: str) -> str:
    """Classify document type using filename, folder, and content signals."""
    filename_lower = path.name.lower()
    folder_lower = path.parent.name.lower()
    text_sample = (text[:3000] if text else "").lower()

    # Content-based (highest priority)
    for pattern, dtype in CONTENT_DOCTYPE_PATTERNS:
        if pattern.search(text[:3000] if text else ""):
            return dtype

    # Filename-based (before folder — specific filenames beat generic folder names)
    for name_key, dtype in FILE_BASED_DOCTYPE.items():
        if name_key in filename_lower:
            return dtype

    # Folder-based fallback
    for folder_key, dtype in FOLDER_DOCTYPE.items():
        if folder_key in folder_lower:
            return dtype

    # Extension-based fallback
    ext = path.suffix.lower()
    if ext == ".msg":
        # Check if it looks internal
        if any(k in filename_lower for k in ("certified mailing", "usps", "return receipt")):
            return "mailroom_record"
        return "internal_communication"

    return "other"


def extract_judge(text: str) -> Optional[str]:
    for pattern, name in JUDGE_PATTERNS:
        if pattern.search(text):
            return name
    return None


def extract_author_recipient(path: Path, text: str) -> Tuple[Optional[str], Optional[str]]:
    """Extract author and recipient from msg/email files or signature blocks."""
    ext = path.suffix.lower()
    author = None
    recipient = None

    if ext == ".msg":
        try:
            import extract_msg as em_lib
            msg = em_lib.Message(str(path))
            author = str(msg.sender or "").strip() or None
            recipient = str(msg.to or "").strip() or None
            return author, recipient
        except Exception:
            pass

    # From/To headers in EML text
    from_m = re.search(r'^From:\s*(.+)$', text, re.MULTILINE)
    to_m = re.search(r'^To:\s*(.+)$', text, re.MULTILINE)
    if from_m:
        author = from_m.group(1).strip()
    if to_m:
        recipient = to_m.group(1).strip()
    if author or recipient:
        return author, recipient

    # Signature block: "BY: Name" or "Sincerely, ..."
    by_m = re.search(r'(?:BY:|Sincerely,)\s*\n\s*(.+)', text, re.MULTILINE)
    if by_m:
        author = by_m.group(1).strip()[:100]

    return author, recipient


def extract_doc_date(text: str, filename: str = "") -> Optional[str]:
    """Extract document date as YYYY-MM-DD from content or filename."""
    # ISO date from filename e.g. 2025-09-11
    m = re.search(r'(\d{4})[-._](\d{2})[-._](\d{2})', filename)
    if m:
        y, mo, d = m.groups()
        if 2020 <= int(y) <= 2030 and 1 <= int(mo) <= 12 and 1 <= int(d) <= 31:
            return f"{y}-{mo}-{d}"

    # Date patterns like "September 12, 2025" or "9/12/2025" or "09.12.25"
    patterns = [
        (re.compile(
            r'(January|February|March|April|May|June|July|August|September|October|November|December)'
            r'\s+(\d{1,2}),?\s+(20\d{2})', re.I
        ), "mname"),
        (re.compile(r'(\d{1,2})/(\d{1,2})/(20\d{2})'), "mdy"),
        (re.compile(r'(20\d{2})[.\-](\d{2})[.\-](\d{2})'), "ymd"),
    ]
    month_map = {
        "january": "01", "february": "02", "march": "03", "april": "04",
        "may": "05", "june": "06", "july": "07", "august": "08",
        "september": "09", "october": "10", "november": "11", "december": "12",
    }

    sample = text[:5000] if text else ""
    for pat, fmt in patterns:
        m = pat.search(sample)
        if m:
            g = m.groups()
            if fmt == "mname":
                mo = month_map.get(g[0].lower(), "00")
                d = g[1].zfill(2)
                y = g[2]
                if mo != "00" and 2020 <= int(y) <= 2030:
                    return f"{y}-{mo}-{d}"
            elif fmt == "mdy":
                mo, d, y = g
                if 2020 <= int(y) <= 2030:
                    return f"{y}-{mo.zfill(2)}-{d.zfill(2)}"
            elif fmt == "ymd":
                y, mo, d = g
                if 2020 <= int(y) <= 2030:
                    return f"{y}-{mo}-{d}"
    return None


def extract_accommodation_type(text: str) -> Optional[str]:
    sample = text[:4000] if text else ""
    found = []
    for atype, pat in ACCOMM_TYPE_PATTERNS.items():
        if pat.search(sample):
            found.append(atype)
    if not found:
        return None
    if len(found) > 1:
        return "multiple"
    return found[0]


def extract_accommodation_action(doc_type: str, text: str) -> Optional[str]:
    if doc_type not in ("accommodation_order", "accommodation_request"):
        return None
    sample = text[:5000] if text else ""
    for action, pat in ACCOMM_ACTION_PATTERNS.items():
        if pat.search(sample):
            return action
    if doc_type == "accommodation_request":
        return "pending"
    return "not_applicable"


def extract_page_count(path: Path) -> Optional[int]:
    if path.suffix.lower() != ".pdf":
        return None
    try:
        result = subprocess.run(
            ["pdfinfo", str(path)],
            capture_output=True, text=True, timeout=15
        )
        m = re.search(r'Pages:\s*(\d+)', result.stdout)
        if m:
            return int(m.group(1))
    except Exception:
        pass
    try:
        import fitz
        doc = fitz.open(str(path))
        n = doc.page_count
        doc.close()
        return n
    except Exception:
        return None


def extract_msg_date(path: Path) -> Optional[str]:
    """Extract date from .msg file."""
    try:
        import extract_msg as em_lib
        msg = em_lib.Message(str(path))
        d = msg.date
        if d:
            m = re.search(r'(\d{4})-(\d{2})-(\d{2})', str(d))
            if m:
                return f"{m.group(1)}-{m.group(2)}-{m.group(3)}"
    except Exception:
        pass
    return None


def build_metadata(
    path: Path,
    text: str,
    doc_type: str,
    extraction_method: str,
    installment: Optional[int],
) -> Dict:
    """Build the metadata JSONB object."""
    now = datetime.now(timezone.utc).isoformat()
    ext = path.suffix.lstrip(".").lower()

    meta: Dict = {
        "source": "samba_dropbox",
        "dropped_at": now,
        "file_type": ext,
        "extraction_method": extraction_method,
        "prr_source": "biia-nelson",
        "prr_request_date": "2025-09-17",
        "prr_scope": "nelson_specific",
    }

    if installment:
        meta["installment_number"] = installment

    # Always include these for this dataset
    meta["biia_docket_number"] = "25-18153"
    meta["biia_claim_number"] = "SX-21824"

    # Date of document
    doc_date = None
    if ext == "msg":
        doc_date = extract_msg_date(path)
    if not doc_date:
        doc_date = extract_doc_date(text, path.name)
    if doc_date:
        meta["date_of_document"] = doc_date

    # Author / recipient
    author, recipient = extract_author_recipient(path, text)
    if author:
        meta["author"] = author[:200]
    if recipient:
        meta["recipient"] = recipient[:200]

    # Judge
    judge = extract_judge(text[:5000] if text else "")
    if judge:
        meta["judge_name"] = judge

    # Accommodation
    accomm_type = extract_accommodation_type(text)
    if accomm_type:
        meta["accommodation_type"] = accomm_type

    accomm_action = extract_accommodation_action(doc_type, text)
    if accomm_action:
        meta["accommodation_action"] = accomm_action

    # Disability type (Nelson records are consistently cognitive_neurological)
    if text and re.search(r'long\s*covid|neurocog|cognitive|neurolog|brain\s+fog|processing\s+speed|executive\s+function', text[:5000], re.I):
        meta["disability_type"] = "cognitive_neurological"
    elif text and re.search(r'disabilit', text[:3000], re.I):
        meta["disability_type"] = "unspecified"

    # Page count
    pc = extract_page_count(path)
    if pc is not None:
        meta["page_count"] = pc

    # No redactions per BIIA statement
    meta["has_redactions"] = False

    return meta


def build_title(path: Path, text: str) -> str:
    """Derive a clean title from content or filename."""
    # Look for obvious document title in content (first ~1000 chars)
    if text:
        top = text[:1500].strip()
        for pat in [
            re.compile(r'^(ORDER\s+[A-Z][A-Z\s/,:\-]+?)(?:\n|$)', re.M),
            re.compile(r'^(LETTER\s+TO\s+[^\n]+?)(?:\n|$)', re.M),
            re.compile(r'^(BIIA\s+[A-Z][A-Z\s]+?)(?:\n|$)', re.M),
        ]:
            m = pat.search(top)
            if m:
                candidate = m.group(1).strip()
                if 10 < len(candidate) < 150:
                    return candidate

    # Fall back to cleaned filename
    stem = path.stem
    # Remove date prefixes like "9.12.25 " or "2025-09-11 "
    stem = re.sub(r'^[\d./-]+\s*', '', stem)
    stem = re.sub(r'[-_]+', ' ', stem)
    stem = re.sub(r'\s+', ' ', stem).strip()
    return stem if stem else path.name


# ──────────────────────────────────────────────────────────────────────────────
# Chunking
# ──────────────────────────────────────────────────────────────────────────────

def chunk_text(text: str, target_chars: int = CHUNK_TARGET_CHARS, overlap: int = CHUNK_OVERLAP_CHARS) -> List[str]:
    """Split text into overlapping chunks respecting paragraph boundaries."""
    if not text:
        return []
    if len(text) <= target_chars:
        return [text]

    paragraphs = re.split(r'\n{2,}', text)
    chunks = []
    current = ""

    for para in paragraphs:
        if not para.strip():
            continue
        candidate = (current + "\n\n" + para).strip() if current else para.strip()
        if len(candidate) <= target_chars:
            current = candidate
        else:
            if current:
                chunks.append(current)
                # Overlap: back up by overlap chars worth of text
                overlap_text = current[-overlap:] if len(current) > overlap else current
                current = overlap_text + "\n\n" + para.strip()
            else:
                # Para itself is too long — split on sentences
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sent in sentences:
                    candidate = (current + " " + sent).strip() if current else sent.strip()
                    if len(candidate) <= target_chars:
                        current = candidate
                    else:
                        if current:
                            chunks.append(current)
                        current = sent.strip()

    if current:
        chunks.append(current)

    return [c for c in chunks if c.strip()]


# ──────────────────────────────────────────────────────────────────────────────
# Embedding client
# ──────────────────────────────────────────────────────────────────────────────

class EmbeddingClient:
    """Local BGE-M3 embedding client (delegates to embedding_service)."""

    def __init__(self):
        self.total_tokens = 0

    async def embed(self, texts: List[str]) -> List[List[float]]:
        if not texts:
            return []
        return embed_texts_sync(texts)

    async def close(self):
        pass  # No HTTP client to close


def vec_literal(embedding: List[float]) -> str:
    return _embedding_vec_literal(embedding)


def count_tokens_approx(text: str) -> int:
    """Rough token count: ~4 chars per token for English text."""
    return max(1, len(text) // 4)


# ──────────────────────────────────────────────────────────────────────────────
# Database
# ──────────────────────────────────────────────────────────────────────────────

async def get_existing_hashes(pool: asyncpg.Pool) -> set:
    """Return set of all content_hashes currently in core.documents."""
    rows = await pool.fetch("SELECT content_hash FROM core.documents WHERE content_hash IS NOT NULL")
    return {r["content_hash"] for r in rows}


async def check_duplicate(pool: asyncpg.Pool, content_hash: str) -> Optional[str]:
    """Return existing document UUID if hash exists, else None."""
    row = await pool.fetchrow(
        "SELECT id::text FROM core.documents WHERE content_hash = $1 LIMIT 1",
        content_hash
    )
    return row["id"] if row else None


async def get_or_create_tag(pool: asyncpg.Pool, name: str, category: str) -> int:
    row = await pool.fetchrow(
        "INSERT INTO core.tags (name, category) VALUES ($1, $2) ON CONFLICT (name) DO UPDATE SET name=EXCLUDED.name RETURNING id",
        name, category
    )
    return row["id"]


async def apply_tag(pool: asyncpg.Pool, doc_id: str, tag_id: int):
    await pool.execute(
        "INSERT INTO core.document_tags (document_id, tag_id) VALUES ($1::uuid, $2) ON CONFLICT DO NOTHING",
        doc_id, tag_id
    )


async def upsert_existing_duplicate(pool: asyncpg.Pool, existing_id: str, prr_tag_id: int, nelson_tag_id: int):
    """Update existing doc to mark it as also in this PRR set."""
    await pool.execute("""
        UPDATE core.documents
        SET metadata = metadata || '{"also_in_prr": true, "prr_source": "biia-nelson"}'::jsonb
        WHERE id = $1::uuid
    """, existing_id)
    await apply_tag(pool, existing_id, prr_tag_id)
    await apply_tag(pool, existing_id, nelson_tag_id)


async def insert_document(pool: asyncpg.Pool, doc: BiiaDocument) -> str:
    """Insert document into core.documents and return its UUID."""
    doc_uuid = uuid.UUID(doc.doc_id)
    await pool.execute("""
        INSERT INTO core.documents
            (id, domain, source_path, filename, document_type, title,
             content_hash, total_chunks, full_content, metadata)
        VALUES (
            $1, $2, $3, $4, $5, $6, $7, $8, $9,
            $10::jsonb
        )
        ON CONFLICT (id) DO NOTHING
    """,
        doc_uuid, "legal", doc.source_path, doc.filename,
        doc.document_type, doc.title, doc.content_hash,
        len(doc.chunks), doc.full_content[:500000],
        json.dumps(doc.metadata),
    )
    return doc.doc_id


async def insert_chunks(pool: asyncpg.Pool, doc_id: str, chunks: List[str], embeddings: List[List[float]], enriched_texts: Optional[List[str]] = None):
    """Insert all chunks with embeddings."""
    doc_uuid = uuid.UUID(doc_id)
    total = len(chunks)
    async with pool.acquire() as conn:
        async with conn.transaction():
            for idx, (chunk_text, emb) in enumerate(zip(chunks, embeddings)):
                chunk_id = hashlib.md5(f"{doc_id}:{idx}".encode()).hexdigest()
                token_count = count_tokens_approx(chunk_text)
                enriched = enriched_texts[idx] if enriched_texts and idx < len(enriched_texts) else None
                await conn.execute("""
                    INSERT INTO core.document_chunks
                        (id, document_id, chunk_index, total_chunks,
                         content, embedded_content, embedding, embedding_model_id, token_count)
                    VALUES (
                        $1, $2, $3, $4, $5, $6,
                        $7::halfvec, $8, $9
                    )
                    ON CONFLICT (id) DO NOTHING
                """,
                    chunk_id, doc_uuid, idx, total,
                    chunk_text, enriched, vec_literal(emb), EMBEDDING_MODEL_ID, token_count,
                )


async def update_total_chunks(pool: asyncpg.Pool, doc_id: str, total: int):
    await pool.execute(
        "UPDATE core.documents SET total_chunks = $1 WHERE id = $2::uuid",
        total, doc_id
    )


# ──────────────────────────────────────────────────────────────────────────────
# Ingestion job tracking
# ──────────────────────────────────────────────────────────────────────────────

async def create_ingestion_job(pool: asyncpg.Pool) -> int:
    row = await pool.fetchrow("""
        INSERT INTO ops.ingestion_jobs (domain, source, status, metadata)
        VALUES ('legal', 'prr_biia_nelson', 'running', $1::jsonb)
        RETURNING id
    """, json.dumps({
        "prr_source": "biia-nelson",
        "prr_scope": "nelson_specific",
        "docket": "25-18153",
        "claim": "SX-21824",
        "date_range": "2025-06-01 to 2025-09-17",
        "installments": 7,
        "source_folder": "09.18.2025 Nelson Request 2518153 Accommodation Records 6.1.25-9.17.25",
    }))
    return row["id"]


async def finish_ingestion_job(pool: asyncpg.Pool, job_id: int, docs_ok: int, chunks_ok: int, docs_failed: int, error_log: List):
    await pool.execute("""
        UPDATE ops.ingestion_jobs
        SET status = $1,
            documents_processed = $2,
            chunks_processed = $3,
            documents_failed = $4,
            completed_at = now(),
            error_log = $5::jsonb
        WHERE id = $6
    """,
        "completed" if docs_failed == 0 else "completed_with_errors",
        docs_ok, chunks_ok, docs_failed,
        json.dumps(error_log),
        job_id,
    )


# ──────────────────────────────────────────────────────────────────────────────
# Tag setup
# ──────────────────────────────────────────────────────────────────────────────

REQUIRED_TAGS = [
    ("biia-prr-nelson",              "prr_source"),
    ("nelson-specific",              "scope"),
    ("accommodation-request",        "document_class"),
    ("accommodation-order",          "document_class"),
    ("accommodation-policy",         "document_class"),
    ("accommodation-denied",         "outcome"),
    ("accommodation-granted",        "outcome"),
    ("accommodation-partial",        "outcome"),
    ("disability-cognitive-neurological", "disability_type"),
    ("disability-physical",          "disability_type"),
    ("disability-sensory",           "disability_type"),
    ("disability-mental-health",     "disability_type"),
    ("disability-developmental",     "disability_type"),
    ("judge-blood",                  "judge"),
    ("judge-dodge",                  "judge"),
    ("judge-watkins",                "judge"),
    ("judge-rosen",                  "judge"),
    ("internal-communication",       "document_class"),
    ("external-communication",       "document_class"),
    ("phone-record",                 "document_class"),
    ("mailroom-record",              "document_class"),
    ("draft-document",               "document_class"),
]

DOCTYPE_TO_TAG = {
    "accommodation_request":  "accommodation-request",
    "accommodation_order":    "accommodation-order",
    "accommodation_policy":   "accommodation-policy",
    "internal_communication": "internal-communication",
    "external_communication": "external-communication",
    "phone_record":           "phone-record",
    "mailroom_record":        "mailroom-record",
    "draft_document":         "draft-document",
}

JUDGE_TO_TAG = {
    "Timothy M. Blood": "judge-blood",
    "Ann Dodge":        "judge-dodge",
    "Brian Watkins":    "judge-watkins",
    "Janice Rosen":     "judge-rosen",
}

ACTION_TO_TAG = {
    "granted":          "accommodation-granted",
    "denied":           "accommodation-denied",
    "partially_granted":"accommodation-partial",
}

DISABILITY_TO_TAG = {
    "cognitive_neurological": "disability-cognitive-neurological",
    "physical_mobility":      "disability-physical",
    "sensory_visual":         "disability-sensory",
    "sensory_hearing":        "disability-sensory",
    "mental_health":          "disability-mental-health",
    "developmental":          "disability-developmental",
}


async def setup_tags(pool: asyncpg.Pool) -> Dict[str, int]:
    """Ensure all required tags exist, return name→id mapping."""
    tag_ids = {}
    for name, category in REQUIRED_TAGS:
        tag_id = await get_or_create_tag(pool, name, category)
        tag_ids[name] = tag_id
    return tag_ids


async def apply_document_tags(pool: asyncpg.Pool, doc: BiiaDocument, tag_ids: Dict[str, int]):
    """Apply all relevant tags to a document."""
    doc_id = doc.doc_id
    meta = doc.metadata

    # Always apply
    await apply_tag(pool, doc_id, tag_ids["biia-prr-nelson"])
    await apply_tag(pool, doc_id, tag_ids["nelson-specific"])

    # Document class
    class_tag = DOCTYPE_TO_TAG.get(doc.document_type)
    if class_tag and class_tag in tag_ids:
        await apply_tag(pool, doc_id, tag_ids[class_tag])

    # Accommodation action (outcome)
    action = meta.get("accommodation_action")
    if action and action in ACTION_TO_TAG:
        tag_name = ACTION_TO_TAG[action]
        if tag_name in tag_ids:
            await apply_tag(pool, doc_id, tag_ids[tag_name])

    # Disability type
    disability = meta.get("disability_type")
    if disability and disability in DISABILITY_TO_TAG:
        tag_name = DISABILITY_TO_TAG[disability]
        if tag_name in tag_ids:
            await apply_tag(pool, doc_id, tag_ids[tag_name])

    # Judge
    judge = meta.get("judge_name")
    if judge:
        judge_tag = JUDGE_TO_TAG.get(judge)
        if judge_tag:
            if judge_tag not in tag_ids:
                # Create on the fly
                tid = await get_or_create_tag(pool, judge_tag, "judge")
                tag_ids[judge_tag] = tid
            await apply_tag(pool, doc_id, tag_ids[judge_tag])
        else:
            # New judge — create tag following pattern judge-lastname
            last = judge.split()[-1].lower()
            new_tag = f"judge-{last}"
            tid = await get_or_create_tag(pool, new_tag, "judge")
            tag_ids[new_tag] = tid
            await apply_tag(pool, doc_id, tid)


# ──────────────────────────────────────────────────────────────────────────────
# File discovery
# ──────────────────────────────────────────────────────────────────────────────

def should_skip(path: Path) -> bool:
    name = path.name
    name_lower = name.lower()
    if name_lower in SKIP_NAMES:
        return True
    if name.startswith("._"):
        return True
    if name.startswith("."):
        return True
    for part in path.parts:
        if part.lower().startswith("__macosx"):
            return True
    return False


def discover_files(root: Path) -> List[Path]:
    """Recursively discover all processable files."""
    files = []
    for f in root.rglob("*"):
        if f.is_file() and not should_skip(f):
            files.append(f)
    return sorted(files)


# ──────────────────────────────────────────────────────────────────────────────
# Validation queries
# ──────────────────────────────────────────────────────────────────────────────

async def run_validation(pool: asyncpg.Pool) -> Dict:
    pattern = "%biia%Nelson%Request%2518153%"

    total = await pool.fetchval(
        "SELECT COUNT(*) FROM core.documents WHERE source_path LIKE $1", pattern
    )

    by_type = await pool.fetch(
        "SELECT document_type, COUNT(*) as cnt FROM core.documents "
        "WHERE source_path LIKE $1 GROUP BY document_type ORDER BY cnt DESC", pattern
    )

    missing_text = await pool.fetch(
        "SELECT filename, source_path FROM core.documents "
        "WHERE source_path LIKE $1 AND (full_content IS NULL OR LENGTH(full_content) < 10)", pattern
    )

    tag_dist = await pool.fetch("""
        SELECT t.name, t.category, COUNT(dt.document_id) as doc_count
        FROM core.tags t
        JOIN core.document_tags dt ON t.id = dt.tag_id
        JOIN core.documents d ON dt.document_id = d.id
        WHERE d.source_path LIKE $1
        GROUP BY t.name, t.category
        ORDER BY t.category, doc_count DESC
    """, pattern)

    action_summary = await pool.fetch("""
        SELECT metadata->>'accommodation_action' as action, COUNT(*) as cnt
        FROM core.documents
        WHERE source_path LIKE $1 AND metadata->>'accommodation_action' IS NOT NULL
        GROUP BY metadata->>'accommodation_action'
    """, pattern)

    judge_dist = await pool.fetch("""
        SELECT metadata->>'judge_name' as judge, COUNT(*) as cnt
        FROM core.documents
        WHERE source_path LIKE $1 AND metadata->>'judge_name' IS NOT NULL
        GROUP BY metadata->>'judge_name'
    """, pattern)

    duplicates = await pool.fetchval("""
        SELECT COUNT(*) FROM core.tags t
        JOIN core.document_tags dt ON t.id = dt.tag_id
        JOIN core.documents d ON dt.document_id = d.id
        WHERE t.name = 'biia-prr-nelson'
        AND d.source_path NOT LIKE $1
        AND d.metadata->>'also_in_prr' = 'true'
    """, pattern)

    return {
        "total": total,
        "by_type": [dict(r) for r in by_type],
        "missing_text": [dict(r) for r in missing_text],
        "tag_dist": [dict(r) for r in tag_dist],
        "action_summary": [dict(r) for r in action_summary],
        "judge_dist": [dict(r) for r in judge_dist],
        "duplicates": duplicates or 0,
    }


# ──────────────────────────────────────────────────────────────────────────────
# Report generation
# ──────────────────────────────────────────────────────────────────────────────

def generate_report(
    stats: Dict,
    validation: Dict,
    failed_files: List[Dict],
    duplicate_files: List[Dict],
    elapsed: float,
) -> str:
    lines = [
        "# BIIA Nelson PRR Ingestion Report",
        "",
        f"**Generated**: {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M:%S UTC')}",
        f"**Source**: BIIA Docket 25-18153 / Claim SX-21824",
        f"**Installments**: 7 (Oct 2025 – Feb 2026)",
        f"**Elapsed**: {elapsed:.1f}s",
        "",
        "---",
        "",
        "## Summary",
        "",
        f"| Metric | Count |",
        f"|--------|-------|",
        f"| Documents processed | {stats['docs_ok']} |",
        f"| Documents failed | {stats['docs_failed']} |",
        f"| Duplicates merged | {stats['docs_dup']} |",
        f"| Total chunks created | {stats['chunks_total']} |",
        f"| Total in DB (this PRR) | {validation.get('total', '?')} |",
        "",
    ]

    lines += [
        "## Document Type Distribution",
        "",
        "| Document Type | Count |",
        "|---------------|-------|",
    ]
    for row in validation.get("by_type", []):
        lines.append(f"| {row['document_type']} | {row['cnt']} |")
    lines.append("")

    lines += [
        "## Accommodation Action Summary",
        "",
        "| Action | Count |",
        "|--------|-------|",
    ]
    for row in validation.get("action_summary", []):
        lines.append(f"| {row['action']} | {row['cnt']} |")
    lines.append("")

    lines += [
        "## Judge Distribution",
        "",
        "| Judge | Documents |",
        "|-------|-----------|",
    ]
    for row in validation.get("judge_dist", []):
        lines.append(f"| {row['judge']} | {row['cnt']} |")
    lines.append("")

    lines += [
        "## Tag Distribution",
        "",
        "| Tag | Category | Doc Count |",
        "|-----|----------|-----------|",
    ]
    for row in validation.get("tag_dist", []):
        lines.append(f"| {row['name']} | {row['category']} | {row['doc_count']} |")
    lines.append("")

    lines += [
        "## Duplicates Found (previously in DB)",
        "",
        f"Documents that already existed in database and received merged tags: **{stats['docs_dup']}**",
        "",
    ]
    if duplicate_files:
        lines += ["| Filename | Existing Document ID |", "|----------|---------------------|"]
        for d in duplicate_files[:50]:
            lines.append(f"| {d['filename']} | {d['existing_id']} |")
    lines.append("")

    lines += [
        "## Extraction Failures (Manual Review Needed)",
        "",
    ]
    if failed_files:
        lines += ["| Filename | Path | Error |", "|----------|------|-------|"]
        for f in failed_files:
            lines.append(f"| {f['filename']} | {f['path']} | {f.get('error', '')} |")
    else:
        lines.append("No extraction failures.")
    lines.append("")

    if validation.get("missing_text"):
        lines += [
            "## Documents with Missing/Short Text in DB",
            "",
            "| Filename | Source Path |",
            "|----------|-------------|",
        ]
        for r in validation["missing_text"]:
            lines.append(f"| {r['filename']} | {r['source_path']} |")
        lines.append("")

    lines += [
        "## Notes",
        "",
        "- BIIA confirmed: no exemptions applied, no redactions made, no pages withheld.",
        "- Disability type for Nelson records: Long COVID / neurocognitive impairment (cognitive_neurological).",
        "- Primary judge: Timothy M. Blood (Assistant Chief IAJ, also ADA coordinator).",
        "- October 9, 2025 BIIA panel order: granted appointment of counsel, ordered email service, denied GAL.",
        "",
    ]

    return "\n".join(lines)


# ──────────────────────────────────────────────────────────────────────────────
# Main pipeline
# ──────────────────────────────────────────────────────────────────────────────

async def main(dry_run: bool = False):
    start_time = time.time()
    log.info("=" * 70)
    log.info("BIIA Nelson PRR Ingestion — Starting")
    log.info(f"Source: {SOURCE_ROOT}")
    log.info(f"Dry run: {dry_run}")
    log.info("=" * 70)

    # Discover files
    all_files = discover_files(SOURCE_ROOT)
    log.info(f"Discovered {len(all_files)} files to process")

    if not all_files:
        log.error("No files found — check source path")
        return

    pool = await asyncpg.create_pool(DATABASE_URL, min_size=2, max_size=8)
    embedder = EmbeddingClient()

    # Setup tags and get ingestion job ID
    tag_ids = await setup_tags(pool)
    log.info(f"Tags ready: {len(tag_ids)} tags")

    job_id = None if dry_run else await create_ingestion_job(pool)
    if job_id:
        log.info(f"Ingestion job ID: {job_id}")

    # Load existing hashes for fast dedup lookup
    existing_hashes = await get_existing_hashes(pool)
    log.info(f"Loaded {len(existing_hashes)} existing content hashes for dedup")

    # Stats
    docs_ok = 0
    docs_failed = 0
    docs_dup = 0
    chunks_total = 0
    failed_files = []
    duplicate_files = []
    error_log = []

    for file_num, fpath in enumerate(all_files, 1):
        log.info(f"[{file_num}/{len(all_files)}] {fpath.name}")

        try:
            # Extract text
            text, method = extract_content(fpath)
            
            if not text or len(text.strip()) < 5:
                log.warning(f"  ⚠ No text extracted from {fpath.name} (method: {method})")
                # Still ingest with empty content (document exists, just no text)
                text = f"[Text extraction failed — {method}]"
                if method == "failed":
                    docs_failed += 1
                    failed_files.append({
                        "filename": fpath.name,
                        "path": str(fpath),
                        "error": "text extraction failed"
                    })
                    error_log.append({"file": str(fpath), "error": "text extraction failed"})
                    continue

            # Compute hash
            content_hash = hashlib.sha256(text.encode("utf-8")).hexdigest()

            # Dedup check
            if content_hash in existing_hashes:
                existing_id = await check_duplicate(pool, content_hash)
                log.info(f"  ↩ Duplicate (hash match) — existing doc {existing_id}")
                docs_dup += 1
                duplicate_files.append({"filename": fpath.name, "existing_id": existing_id})
                if not dry_run and existing_id:
                    await upsert_existing_duplicate(pool, existing_id, tag_ids["biia-prr-nelson"], tag_ids["nelson-specific"])
                continue

            # Detect installment
            installment = detect_installment(fpath)

            # Classify
            doc_type = classify_document(fpath, text)

            # Build metadata
            meta = build_metadata(fpath, text, doc_type, method, installment)

            # Build title
            title = build_title(fpath, text)

            # Chunk
            chunks = chunk_text(text)
            if not chunks:
                chunks = [text[:2500]]

            # Build document
            doc_id = str(uuid.uuid4())
            doc = BiiaDocument(
                doc_id=doc_id,
                source_path=str(fpath),
                filename=fpath.name,
                document_type=doc_type,
                title=title,
                content_hash=content_hash,
                full_content=text,
                metadata=meta,
                chunks=chunks,
                extraction_method=method,
            )

            log.info(f"  → {doc_type} | {len(chunks)} chunks | {method} | installment {installment}")

            if dry_run:
                log.info(f"  [DRY RUN] Would insert: {doc.title[:60]}")
                docs_ok += 1
                chunks_total += len(chunks)
                existing_hashes.add(content_hash)
                continue

            # Contextual Retrieval: generate context and enrich chunks
            context = generate_context_sync(
                title=title,
                domain="legal",
                document_type=doc_type,
                content_preview=text[:3000],
                case_number=None,
            )
            enriched_texts = enrich_chunks(context, chunks)

            # Generate embeddings in batches (embed enriched texts)
            all_embeddings = []
            for batch_start in range(0, len(enriched_texts), EMBEDDING_BATCH):
                batch = enriched_texts[batch_start:batch_start + EMBEDDING_BATCH]
                try:
                    embs = await embedder.embed(batch)
                    all_embeddings.extend(embs)
                except Exception as e:
                    log.error(f"  Embedding failed for batch {batch_start}: {e}")
                    # Fill with zeros for failed batches — better to have doc without embeddings
                    all_embeddings.extend([[0.0] * EMBEDDING_DIMS] * len(batch))

            # Insert document
            await insert_document(pool, doc)

            # Insert chunks (with enriched texts for embedded_content)
            await insert_chunks(pool, doc_id, chunks, all_embeddings, enriched_texts=enriched_texts)

            # Update total_chunks
            await update_total_chunks(pool, doc_id, len(chunks))

            # Apply tags
            await apply_document_tags(pool, doc, tag_ids)

            docs_ok += 1
            chunks_total += len(chunks)
            existing_hashes.add(content_hash)
            log.info(f"  ✓ Inserted doc {doc_id[:8]}… with {len(chunks)} chunks")

        except Exception as e:
            log.exception(f"  ✗ Failed to process {fpath.name}: {e}")
            docs_failed += 1
            failed_files.append({"filename": fpath.name, "path": str(fpath), "error": str(e)})
            error_log.append({"file": str(fpath), "error": str(e)})

    # Finish ingestion job
    if not dry_run and job_id:
        await finish_ingestion_job(pool, job_id, docs_ok, chunks_total, docs_failed, error_log[:100])
        log.info(f"Ingestion job {job_id} marked completed")

    elapsed = time.time() - start_time

    # Stats summary
    stats = {
        "docs_ok": docs_ok,
        "docs_failed": docs_failed,
        "docs_dup": docs_dup,
        "chunks_total": chunks_total,
    }

    log.info("=" * 70)
    log.info(f"INGESTION COMPLETE in {elapsed:.1f}s")
    log.info(f"  Inserted:   {docs_ok} documents")
    log.info(f"  Duplicates: {docs_dup} merged")
    log.info(f"  Failed:     {docs_failed}")
    log.info(f"  Chunks:     {chunks_total}")
    log.info(f"  Tokens:     {embedder.total_tokens:,}")
    log.info("=" * 70)

    # Validation
    log.info("Running validation queries…")
    if not dry_run:
        validation = await run_validation(pool)
        log.info(f"Validation: {validation['total']} docs in DB from this PRR")
        log.info(f"  By type: {[(r['document_type'], r['cnt']) for r in validation['by_type']]}")
    else:
        validation = {}

    # Generate report
    report = generate_report(stats, validation, failed_files, duplicate_files, elapsed)
    if not dry_run:
        REPORT_PATH.parent.mkdir(parents=True, exist_ok=True)
        REPORT_PATH.write_text(report, encoding="utf-8")
        log.info(f"Report saved to: {REPORT_PATH}")
    else:
        log.info("=== DRY RUN REPORT ===")
        print(report)

    await embedder.close()
    await pool.close()


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Ingest BIIA Nelson PRR records")
    parser.add_argument("--dry-run", action="store_true", help="Extract and classify without writing to DB")
    args = parser.parse_args()

    # Load .env
    env_path = Path("/opt/wdws/.env")
    if env_path.exists():
        for line in env_path.read_text().splitlines():
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, _, v = line.partition("=")
                os.environ.setdefault(k.strip(), v.strip())

    asyncio.run(main(dry_run=args.dry_run))
